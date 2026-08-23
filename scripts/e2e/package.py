"""What PowerPoint sees.

Ptium writes the .pptx itself, so the only honest check is to read one back with
a library that is not ours. python-pptx parses the same OOXML PowerPoint does: if
it finds a real table, real notes and real placeholders, so will PowerPoint — and
anything malformed enough to make PowerPoint offer to repair the file tends to
make python-pptx raise here first.

    pip install python-pptx
    python3 scripts/e2e/package.py
"""
import html, io, json, os, re, sys, time, urllib.request, urllib.error, uuid, zipfile

BASE = os.environ.get("PTIUM_URL", "http://localhost:8099").rstrip("/") + "/api/v1"
SECRET = os.environ.get("PTIUM_DEV_SECRET", "devsecret-devsecret-devsecret-devsecret")
RUN = str(int(time.time()))[-6:]
failures = []

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx is not installed: pip install python-pptx")
    sys.exit(0)


def call(method, path, body=None, raw=False):
    data = json.dumps(body).encode() if body is not None else None
    headers = {"X-Ptium-Dev-Secret": SECRET}
    if data:
        headers["Content-Type"] = "application/json"
    request = urllib.request.Request(BASE + path, data=data, headers=headers, method=method)
    with urllib.request.urlopen(request) as response:
        payload = response.read()
    if raw:
        return payload
    return json.loads(payload)["data"] if payload else None


source = """# 패키지 점검 {run}
@cover
> 내보낸 파일을 다른 도구로 읽습니다

# 분기 실적
::table 분기 실적
- 항목 | 1분기 | 2분기
- 매출 | 120억 | 140억
- 영업이익 | 12억 | 18억
::
!source 1 | 2026 시장 조사 보고서 | p.42
!notes 표는 편집 가능한 표여야 합니다.

# 핵심 지표
::kpi 한눈에
- 전환 대상 | 42개
- 예상 절감 | 18%
::
- 요점 하나

# 분기 매출
::columns 분기 매출
- 1분기 | 1180
- 2분기 | 1240
- 3분기 | 1390
::

# 마무리
@closing
> 감사합니다
""".format(run=RUN)

deck = call("POST", "/presentations", {"title": f"패키지 점검 {RUN}", "prompt": "패키지",
                                       "requestedSlideCount": 5, "language": "ko"})
call("PUT", f"/presentations/{deck['id']}/source", {"source": source})
package = call("GET", f"/presentations/{deck['id']}/export?format=pptx", raw=True)
path = os.path.join(os.environ.get("TMPDIR", "/tmp"), f"ptium-package-{RUN}.pptx")
with open(path, "wb") as handle:
    handle.write(package)
print(f"exported {len(package):,} bytes → {path}")

try:
    read = Presentation(path)
except Exception as error:
    print(f" ✗ python-pptx cannot open the export: {error}")
    sys.exit(1)

print(f"slides: {len(read.slides)}  size: {read.slide_width}x{read.slide_height}")
if len(read.slides) != 5:
    failures.append(f"the package holds {len(read.slides)} slides, expected 5")

tables, charts, numbered, noted, titles = 0, 0, 0, 0, []
undescribed = []  # objects a screen reader would find empty


def alt_text(shape):
    marker = "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
    element = shape._element.find(".//" + marker)
    return (element.get("descr") or "").strip() if element is not None else ""
cover_numbered = False
for index, slide in enumerate(read.slides, 1):
    for shape in slide.shapes:
        if shape.has_table:
            tables += 1
            table = shape.table
            cells = [[cell.text for cell in row.cells] for row in table.rows]
            print(f"  slide{index}: table {len(table.rows)}x{len(table.columns)} {cells[0]}")
            if cells[0] != ["항목", "1분기", "2분기"]:
                failures.append(f"the table's header reads {cells[0]}")
            if len(table.rows) != 3:
                failures.append(f"the table has {len(table.rows)} rows, expected 3")
            # A grouped table is a table nobody can click into.
            if shape.shape_type == 6:
                failures.append("the table is inside a group")
        # A chart someone cannot change the numbers of is a picture of a chart.
        if shape.has_chart:
            charts += 1
            chart = shape.chart
            plot = chart.plots[0]
            values = [list(series.values) for series in plot.series]
            print(f"  slide{index}: chart {chart.chart_type} {list(plot.categories)} {values}")
            if list(plot.categories) != ["1분기", "2분기", "3분기"]:
                failures.append(f"the chart's categories read {list(plot.categories)}")
            if values != [[1180.0, 1240.0, 1390.0]]:
                failures.append(f"the chart plots {values}")
            if not chart.part.part_related_by(
                    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/package"):
                failures.append("the chart carries no workbook, so nobody can edit its data")
        # PowerPoint's accessibility check reports a drawing with no alternative
        # text as an error, which keeps a deck out of a public-sector filing.
        if not shape.is_placeholder and not alt_text(shape):
            undescribed.append(f"slide{index}: {shape.shape_type}")
        if shape.is_placeholder and shape.placeholder_format.type is not None:
            if str(shape.placeholder_format.type).startswith("SLIDE_NUMBER"):
                numbered += 1
                if index == 1:
                    cover_numbered = True
        if shape.has_text_frame and shape.text_frame.text.strip() and not titles:
            pass
    if slide.shapes.title is not None and slide.shapes.title.text.strip():
        titles.append(slide.shapes.title.text.strip())
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
        noted += 1
        # A source travels with the claim into the notes pane, where a presenter
        # looks when someone asks where a number came from.
        if index == 2 and "2026 시장 조사 보고서" not in slide.notes_slide.notes_text_frame.text:
            failures.append("the slide's source is not on its notes page")

print(f"titles: {titles}")
print(f"tables: {tables}  charts: {charts}  numbered slides: {numbered}  slides with notes: {noted}")
print(f"objects without alt text: {len(undescribed)}")
for entry in undescribed:
    failures.append(f"{entry} carries no alternative text")
if tables != 1:
    failures.append(f"found {tables} real tables, expected 1")
if charts != 1:
    failures.append(f"found {charts} real charts, expected 1")
# The cover never carries one, and a layout that opts out of the master's shapes
# — the shipped closing page paints its own full-bleed background — has no place
# for one either. What must hold is that ordinary content slides are numbered and
# the cover is not.
if numbered < 1:
    failures.append("no slide carries a page number")
if cover_numbered:
    failures.append("the cover carries a page number")
if noted < 1:
    failures.append("no slide carries speaker notes")
if len(titles) != 5:
    failures.append(f"read {len(titles)} slide titles, expected 5")

# The preview is what the author judges the deck by and the file is what the
# room sees. Every fix of the last week was checked in the preview alone, so
# the two are compared here line by line: a component drawn in one and missing
# from the other is a deck that changes on the way out.
print("── the preview and the file say the same thing ──")


def texts_in(shapes):
    """Every line of text a slide draws, groups included."""
    found = []
    for shape in shapes:
        if getattr(shape, "shape_type", None) is not None and str(shape.shape_type).startswith("GROUP"):
            found.extend(texts_in(shape.shapes))
            continue
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                line = "".join(run.text for run in paragraph.runs).strip()
                if line:
                    found.append(line)
        if getattr(shape, "has_chart", False) and shape.has_chart:
            # A real chart keeps its labels in the chart part rather than in
            # shape text, which is the whole point of it being a real chart.
            for plot in shape.chart.plots:
                found.extend(str(category) for category in plot.categories)
                for series in plot.series:
                    found.extend(format(value, "g") for value in series.values if value is not None)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        found.append(cell.text.strip())
    return found


def preview_lines(presentation_id, position):
    svg = call("GET", f"/presentations/{presentation_id}/preview.svg?slide={position}&width=1200", raw=True).decode()
    lines = []
    for block in re.findall(r"<text[^>]*>(.*?)</text>", svg, re.S):
        for piece in re.split(r"</?tspan[^>]*>", block):
            piece = re.sub(r"<[^>]+>", "", piece).strip()
            if piece:
                lines.append(html.unescape(piece))
    return lines


for index, slide in enumerate(read.slides, 1):
    in_file = {line.replace(" ", "") for line in texts_in(slide.shapes)}
    # The preview draws the bullet mark; the file leaves it to PowerPoint's own
    # list formatting, which is the same line either way.
    drawn_clean = lambda line: line.lstrip("•▪-–— ").replace(" ", "")
    drawn = preview_lines(deck["id"], index)
    missing = [line for line in drawn
               if not any(drawn_clean(line) in packed for packed in in_file)]
    # The page number is drawn by the preview and written as a placeholder the
    # renderer fills, so it is not compared.
    missing = [line for line in missing if line.strip() != str(index)]
    if missing:
        failures.append(f"slide {index} draws {missing[:4]} in the preview and not in the file")
print(f"   compared {len(read.slides)} slides")

# A deck reaches Ptium three ways: written by the model, typed by hand, and
# carried in from a file. The third is the one nobody watches, so it is walked
# end to end here: a file built by another tool goes in, comes out, has the
# source Ptium hides inside it stripped away, and goes in again — which is the
# path a deck takes when someone edits the exported file in PowerPoint first.
print("── a deck carried in from another tool ──")
try:
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError:
    print("   python-pptx is missing its chart helpers; skipped")
else:
    foreign = Presentation()
    cover = foreign.slides.add_slide(foreign.slide_layouts[0])
    cover.shapes.title.text = f"외부 도구로 만든 덱 {RUN}"
    cover.placeholders[1].text = "가져오기 점검용"
    listed = foreign.slides.add_slide(foreign.slide_layouts[1])
    listed.shapes.title.text = "현황 요약"
    frame = listed.placeholders[1].text_frame
    frame.text = "검색 실패율 31%"
    for line in ["문서 12,400건", "이관 4개월 소요"]:
        frame.add_paragraph().text = line
    listed.notes_slide.notes_text_frame.text = "여기서 실패율을 강조합니다."
    tabled = foreign.slides.add_slide(foreign.slide_layouts[5])
    tabled.shapes.title.text = "대안 비교"
    grid = tabled.shapes.add_table(2, 3, Inches(0.5), Inches(2), Inches(9), Inches(1.5)).table
    for row, line in enumerate([["항목", "현행 유지", "외부 SaaS"], ["초기 비용", "0원", "1억 5천만 원"]]):
        for column, cell in enumerate(line):
            grid.cell(row, column).text = cell
    charted = foreign.slides.add_slide(foreign.slide_layouts[5])
    charted.shapes.title.text = "분기 검색량"
    numbers = CategoryChartData()
    numbers.categories = ["1분기", "2분기", "3분기"]
    numbers.add_series("검색량", (1180, 1240, 1390))
    charted.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(8), Inches(4), numbers)
    carried = os.path.join(os.environ.get("TMPDIR", "/tmp"), f"ptium-foreign-{RUN}.pptx")
    foreign.save(carried)

    def import_file(path, name):
        boundary = uuid.uuid4().hex
        payload = (f"--{boundary}\r\nContent-Disposition: form-data; name=\"file\"; filename=\"{name}\"\r\n"
                   "Content-Type: application/vnd.openxmlformats-officedocument.presentationml.presentation\r\n\r\n").encode()
        payload += open(path, "rb").read() + f"\r\n--{boundary}--\r\n".encode()
        request = urllib.request.Request(BASE + "/presentations/import", data=payload,
                                         headers={"X-Ptium-Dev-Secret": SECRET,
                                                  "Content-Type": f"multipart/form-data; boundary={boundary}"},
                                         method="POST")
        with urllib.request.urlopen(request) as response:
            return json.loads(response.read())["data"]

    brought = import_file(carried, "foreign.pptx")
    deck_id = (brought.get("presentation") or {}).get("id", "")
    source = (call("GET", f"/presentations/{deck_id}/source") or {}).get("source", "")
    for wanted in ["# 현황 요약", "- 검색 실패율 31%", "!notes 여기서 실패율을", "::table", "- 초기 비용 | 0원 | 1억 5천만 원", "::columns", "- 2분기 | 1240"]:
        if wanted not in source:
            failures.append(f"the deck carried in from another tool lost {wanted!r}")
    print(f"   {brought.get('slides')} slides in, table and chart kept")

    # Out and back in again, with the source Ptium hides in the file removed.
    exported = call("GET", f"/presentations/{deck_id}/export?format=pptx", raw=True)
    stripped_path = os.path.join(os.environ.get("TMPDIR", "/tmp"), f"ptium-stripped-{RUN}.pptx")
    inside = zipfile.ZipFile(io.BytesIO(exported))
    with zipfile.ZipFile(stripped_path, "w", zipfile.ZIP_DEFLATED) as out:
        for item in inside.infolist():
            if item.filename == "ppt/ptiumSource.xml":
                continue
            data = inside.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(b'<Override PartName="/ppt/ptiumSource.xml" ContentType="application/xml"/>', b"")
            out.writestr(item, data)
    again = import_file(stripped_path, "stripped.pptx")
    reread = (call("GET", f"/presentations/{(again.get('presentation') or {}).get('id', '')}/source") or {}).get("source", "")
    for wanted in ["- 검색 실패율 31%", "!notes 여기서 실패율을", "- 초기 비용 | 0원 | 1억 5천만 원", "- 2분기 | 1240"]:
        if wanted not in reread:
            failures.append(f"reading Ptium's own drawing back lost {wanted!r}")
    print("   read back from the drawing alone, nothing lost")

print()
print(f"{len(failures)} failures")
for failure in failures:
    print(" ✗ " + failure)
sys.exit(1 if failures else 0)
