"""An end-to-end sweep of the REST surface, run against a real server.

Every check states what it expects. A failure prints the request, the status and
the body, so the next step is reading code rather than reproducing.
"""
import csv, json, os, sys, urllib.parse, urllib.request, urllib.error, zipfile, io, zlib, struct, time, re

BASE = os.environ.get("PTIUM_URL", "http://localhost:8099").rstrip("/") + "/api/v1"
RUN = str(int(time.time()))[-6:]  # each run works on its own names
SECRET = os.environ.get("PTIUM_DEV_SECRET", "devsecret-devsecret-devsecret-devsecret")
failures, checks = [], 0

def call(method, path, body=None, raw=False, files=None, expect=None, note="", headers=None):
    """headers speaks as somebody else: most checks are the development
    identity, and a few need another account or none at all."""
    global checks
    url = path if path.startswith("http") else BASE + path
    headers = dict(headers) if headers else {"X-Ptium-Dev-Secret": SECRET}
    data = None
    if files:
        boundary = "----ptium-e2e"
        parts = []
        for name, (filename, content, ctype) in files.items():
            parts.append(f"--{boundary}\r\nContent-Disposition: form-data; name=\"{name}\"".encode())
            if filename:
                parts[-1] += f"; filename=\"{filename}\"\r\nContent-Type: {ctype}\r\n\r\n".encode()
            else:
                parts[-1] += b"\r\n\r\n"
            parts[-1] += content if isinstance(content, bytes) else content.encode()
            parts[-1] += b"\r\n"
        parts.append(f"--{boundary}--\r\n".encode())
        data = b"".join(parts)
        headers["Content-Type"] = f"multipart/form-data; boundary={boundary}"
    elif body is not None:
        data = json.dumps(body).encode()
        headers["Content-Type"] = "application/json"
    request = urllib.request.Request(url, data=data, headers=headers, method=method)
    try:
        with urllib.request.urlopen(request) as response:
            status, payload = response.status, response.read()
    except urllib.error.HTTPError as error:
        status, payload = error.code, error.read()
    checks += 1
    if expect is not None and status not in (expect if isinstance(expect, (list, tuple)) else [expect]):
        failures.append(f"{method} {path} → {status}, want {expect}. {note}\n    {payload[:400]!r}")
    if raw:
        return status, payload
    try:
        return status, json.loads(payload) if payload else None
    except Exception:
        return status, payload

def data_of(result):
    status, body = result
    return (body or {}).get("data") if isinstance(body, dict) else None

def png(rgb=(200, 120, 60), size=8):
    def chunk(kind, payload):
        body = kind + payload
        return struct.pack(">I", len(payload)) + body + struct.pack(">I", zlib.crc32(body) & 0xffffffff)
    rows = b"".join(b"\x00" + bytes(rgb) * size for _ in range(size))
    # Uploading identical bytes is a reuse, by design. Two runs a while apart
    # used to produce the same picture — the colour came from three digits of the
    # clock — and the later run then worked on the earlier run's image, under the
    # earlier run's name. The run's own text goes in a comment so every run
    # uploads a picture no other run has.
    return (b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", struct.pack(">IIBBBBB", size, size, 8, 2, 0, 0, 0))
            + chunk(b"tEXt", b"ptium-run\x00" + RUN.encode())
            + chunk(b"IDAT", zlib.compress(rows)) + chunk(b"IEND", b""))

def pdf_of(pages, image_only=False):
    """A PDF whose pages say the given lines, written the way the generators
    people actually use write Korean: one byte for a character that fits in one,
    two for a character that does not."""
    def shown(line):
        return "".join(f"{ord(c):02x}" if ord(c) < 0x80 else f"{ord(c):04x}" for c in line)
    objects = ["<</Type /Catalog /Pages 2 0 R>>"]
    kids = " ".join(f"{3 + index * 2} 0 R" for index in range(len(pages)))
    objects.append(f"<</Type /Pages /Kids [{kids}] /Count {len(pages)}>>")
    font = 3 + len(pages) * 2
    for index, lines in enumerate(pages):
        if image_only:
            content = "q 612 0 0 792 0 0 cm /Im1 Do Q"
        else:
            content = "".join(
                f"BT /F1 12 Tf 1 0 0 1 72 {700 - row * 20} Tm <{shown(line)}> Tj ET\n"
                for row, line in enumerate(lines))
        objects.append(f"<</Type /Page /Parent 2 0 R /Resources <</Font <</F1 {font} 0 R>>>> "
                       f"/Contents {4 + index * 2} 0 R>>")
        objects.append(f"<</Length {len(content)}>>\nstream\n{content}\nendstream")
    objects.append("<</Type /Font /Subtype /Type1 /BaseFont /Helvetica>>")
    out = "%PDF-1.7\n"
    for index, body in enumerate(objects):
        out += f"{index + 1} 0 obj\n{body}\nendobj\n"
    return (out + "trailer<</Root 1 0 R>>\n%%EOF\n").encode()


print("── identity and settings ──")
me = data_of(call("GET", "/me", expect=200))
print("   signed in as", (me or {}).get("email"), "admin:", (me or {}).get("isAdmin"))
call("GET", "/profile", expect=200)
call("PATCH", "/profile", {"displayName": "E2E Tester", "company": "Ptium", "jobTitle": "QA"}, expect=200)
public_settings = data_of(call("GET", "/settings", expect=200)) or {}
# The personalisation screen shows a colour in a swatch whether or not anybody
# picked one — a deployment that never opened the branding screen still carries
# the colour this product seeds. The drawing ignores exactly that colour, so a
# screen that cannot tell which value it is goes back to promising a colour that
# will never be drawn.
if not re.fullmatch(r"#[0-9A-Fa-f]{6}", str(public_settings.get("branding.seeded_brand_color", ""))):
    failures.append(f"settings do not say which colour nobody chose: {public_settings.get('branding.seeded_brand_color')!r}")
call("GET", "/admin/settings", expect=[200, 403])
# One page a site with no internet can hand to somebody who cannot see it. It
# must carry what is wrong and never a secret.
status, report_body = call("GET", "/admin/report", expect=[200, 403])
report = data_of((status, report_body)) or {}
if report:
    for field in ("version", "schema", "overview", "usage", "tidy", "changedSettings", "designs", "modelHost"):
        if field not in report:
            failures.append(f"the deployment report does not carry {field}")
    schema = report.get("schema") or {}
    if (schema.get("applied") or 0) < 1:
        failures.append("the report says no migrations have been applied")
    _, as_text = call("GET", "/admin/report?format=md", raw=True, expect=200)
    text = as_text.decode("utf-8", "replace") if isinstance(as_text, (bytes, bytearray)) else str(as_text)
    if "Ptium 배포 점검" not in text:
        failures.append("the report as a file does not read as a report")
    if "비밀 값이 들어 있지 않습니다" not in text:
        failures.append("the report does not say that it carries no secret")
    for leaked in (SECRET, "postgres://", "sk-sweep-secret"):
        if leaked and leaked in text:
            failures.append("the deployment report carries a secret")
    for named in ("unusedImages", "trashed", "failedOldDecks"):
        if named + ":" in text:
            failures.append(f"the report calls a pile by its field name: {named}")

# What has accumulated and is going nowhere. This answer deletes nothing; it is
# what a decision about keeping things would be made from.
tidy = data_of(call("GET", "/admin/tidy", expect=[200, 403])) or {}
if tidy:
    kinds = {row.get("kind"): row for row in (tidy.get("items") or [])}
    for wanted in ("trashed", "failedOldDecks", "untouchedDrafts", "expiredLinks",
                   "unusedImages", "unusedImagesOverAMonth"):
        if wanted not in kinds:
            failures.append(f"the tidy preview does not count {wanted}")
    for kind, row in kinds.items():
        if (row.get("count") or 0) < 0 or (row.get("bytes") or 0) < 0:
            failures.append(f"{kind} counts {row.get('count')!r} / {row.get('bytes')!r}")
        if (row.get("count") or 0) == 0 and row.get("oldest"):
            failures.append(f"{kind} has nothing in it and still names an oldest one")
    older = kinds.get("unusedImagesOverAMonth") or {}
    all_unused = kinds.get("unusedImages") or {}
    if (older.get("count") or 0) > (all_unused.get("count") or 0):
        failures.append("more images are unused for a month than are unused at all")
    # It reads and never removes: asking twice leaves the same numbers.
    again = data_of(call("GET", "/admin/tidy", expect=200)) or {}
    if {row.get("kind"): row.get("count") for row in (again.get("items") or [])} != \
       {row.get("kind"): row.get("count") for row in (tidy.get("items") or [])}:
        failures.append("reading the tidy preview changed what it counts")

# The designs this deployment writes decks in. An upload is private to whoever
# uploaded it and the standard was a text field whose value meant nothing to
# anyone reading it, so an operator could see neither.
designs = data_of(call("GET", "/admin/templates?limit=100", expect=[200, 403])) or []
if designs:
    builtin = [one for one in designs if one.get("kind") == "builtin"]
    if len(builtin) < 10:
        failures.append(f"the deployment's designs list holds {len(builtin)} built-in designs")
    for one in designs:
        if one.get("decks") is None or one.get("recent") is None:
            failures.append(f"{one.get('name')!r} does not say how much it is used")
        if (one.get("recent") or 0) > (one.get("decks") or 0):
            failures.append(f"{one.get('name')!r} was used more in thirty days than ever")
    standard_now = [one for one in designs if one.get("standard")]
    if len(standard_now) > 1:
        failures.append(f"{len(standard_now)} designs claim to be the standard")
    call("GET", "/admin/templates?kind=nonsense", expect=422)
    # A built-in design belongs to everybody already, and says so rather than
    # pretending to change.
    if builtin:
        call("POST", f"/admin/templates/{builtin[0]['id']}/shared", {"shared": True}, expect=422)
        # Making one the standard is the same setting the settings screen holds,
        # and the trail says which design it became.
        kept = {row["key"]: row.get("value") for row in (data_of(call("GET", "/admin/settings", expect=200)) or [])}
        call("POST", f"/admin/templates/{builtin[0]['id']}/standard", {}, expect=200)
        after = data_of(call("GET", "/admin/templates?kind=builtin&limit=100", expect=200)) or []
        if not any(one.get("id") == builtin[0]["id"] and one.get("standard") for one in after):
            failures.append("a design just made the standard does not read as the standard")
        call("PUT", "/admin/settings",
             {"values": {"generation.default_theme": kept.get("generation.default_theme", "slate-classic")}}, expect=200)

# What this deployment has been doing. The overview says what is true now; a
# week of decks — how many, how many failed, how long they took, who asked — was
# nowhere, and on a self-hosted model that time is what running it costs.
usage = data_of(call("GET", "/admin/usage?days=7", expect=[200, 403])) or {}
if usage:
    days = usage.get("days") or []
    if len(days) != 7:
        failures.append(f"seven days of usage came back as {len(days)} rows")
    if sum(day.get("generated", 0) for day in days) != usage.get("generated"):
        failures.append("the day rows and the total do not agree")
    for day in days:
        if day.get("failed", 0) > day.get("generated", 0):
            failures.append(f"{day.get('day')} reports more failures than decks")
        if day.get("slowestSeconds", 0) < day.get("medianSeconds", 0):
            failures.append(f"{day.get('day')} says its slowest deck was faster than its middle one")
    if usage.get("timed", 0) > usage.get("generated", 0):
        failures.append("more generations recorded a duration than exist")
    for group in ("owners", "designs", "failures"):
        for row in (usage.get(group) or []):
            if not str(row.get("name", "")).strip():
                failures.append(f"a row in {group} has nothing to call it")
    call("GET", "/admin/usage?days=0", expect=422)
    call("GET", "/admin/usage?days=999", expect=422)

# Every link this deployment has handed out. Only the deck's owner could see
# their own, so nobody could answer what is readable outside — or close a link
# left open by somebody who has gone.
watched = data_of(call("GET", "/presentations?limit=5", expect=200)) or []
with_slides = [row for row in watched if (row.get("slideCount") or 0) > 0]
if with_slides:
    deck_for_link = with_slides[0]
    made = data_of(call("POST", f"/presentations/{deck_for_link['id']}/shares",
                        {"label": f"감독 점검 {RUN}"}, expect=201)) or {}
    listed = data_of(call("GET", "/admin/shares?state=open&limit=100", expect=200)) or []
    mine = [one for one in listed if one.get("id") == made.get("id")]
    if not mine:
        failures.append("a link somebody just made is not in the deployment's list")
    else:
        one = mine[0]
        if one.get("state") != "open":
            failures.append(f"a link with no day on it reads as {one.get('state')!r}")
        if not (one.get("ownerEmail") or one.get("ownerName")):
            failures.append("a link does not say who made it")
        if one.get("deckTitle") != deck_for_link.get("title"):
            failures.append("a link does not say which deck it opens")
    # Closing it is the operator's, whoever made it, and it says whether it did.
    answered = data_of(call("POST", f"/admin/shares/{made.get('id')}/close", {}, expect=200)) or {}
    if not answered.get("closed"):
        failures.append("closing an open link reported that it closed nothing")
    again = data_of(call("POST", f"/admin/shares/{made.get('id')}/close", {}, expect=200)) or {}
    if again.get("closed"):
        failures.append("closing an already closed link reported that it closed it")
    closed_list = data_of(call("GET", "/admin/shares?state=revoked&limit=100", expect=200)) or []
    if not any(one.get("id") == made.get("id") for one in closed_list):
        failures.append("a closed link is not among the closed ones")
    call("GET", "/admin/shares?state=nonsense", expect=422)
# What a settings change was. The trail used to record "settings.update_batch
# count=1": which setting, what it had been and what it became were nowhere, and
# these settings decide how every deck in the deployment is written.
current = {row["key"]: row.get("value") for row in (data_of(call("GET", "/admin/settings", expect=[200, 403])) or [])}
if "ai.timeout_seconds" in current:
    kept = current["ai.timeout_seconds"]
    call("PUT", "/admin/settings", {"values": {"ai.timeout_seconds": 600 if kept != 600 else 900}}, expect=200)
    changes = data_of(call("GET", "/admin/settings/changes?limit=5", expect=200)) or []
    mine = [c for c in changes if (c.get("metadata") or {}).get("key") == "ai.timeout_seconds"]
    if not mine:
        failures.append("a settings change is not in the settings' own trail")
    else:
        detail = mine[0].get("metadata") or {}
        if detail.get("from") != kept:
            failures.append(f"the trail says the timeout was {detail.get('from')!r}, it was {kept!r}")
        if not mine[0].get("actorName") and not mine[0].get("actorEmail"):
            failures.append("a settings change does not say who made it")
        # And putting it back is one press, not retyping the old value.
        call("POST", f"/admin/settings/changes/{mine[0]['id']}/revert", {}, expect=200)
        after = {row["key"]: row.get("value") for row in (data_of(call("GET", "/admin/settings", expect=200)) or [])}
        if after.get("ai.timeout_seconds") != kept:
            failures.append(f"reverting left the timeout at {after.get('ai.timeout_seconds')!r}, not {kept!r}")
    # A secret's value is never written down, so it cannot be put back.
    call("PUT", "/admin/settings", {"values": {"ai.api_key": "sk-sweep-secret"}}, expect=200)
    secret = [c for c in (data_of(call("GET", "/admin/settings/changes?limit=5", expect=200)) or [])
              if (c.get("metadata") or {}).get("key") == "ai.api_key"]
    if not secret:
        failures.append("a secret changing is not in the trail at all")
    else:
        if "sk-sweep-secret" in json.dumps(secret[0], ensure_ascii=False):
            failures.append("a secret's value was written into the audit trail")
        call("POST", f"/admin/settings/changes/{secret[0]['id']}/revert", {}, expect=422)
    call("PUT", "/admin/settings", {"values": {"ai.api_key": " "}}, expect=200)
# The overview used to call a deck being written "waiting", so a healthy
# half-hour generation read as a stall an operator was told to go and check.
overview = data_of(call("GET", "/admin/overview", expect=[200, 403])) or {}
if overview:
    for field in ("oldestQueuedSeconds", "quietestGenerationSeconds"):
        if not isinstance(overview.get(field), int):
            failures.append(f"the overview does not answer {field}: {overview.get(field)!r}")
# A setting this deployment will not honour used to be stored and shown back:
# an administrator could read "생성 후 자동 수정 500" off their own screen while
# generation ran three passes. Each of these is refused now, and what a
# deployment does honour still saves.
for key, refused, honoured in [
    ("ai.timeout_seconds", 99999, 300),
    ("ai.max_output_tokens", 90, 8000),
    ("ai.reasoning", "thinking-hard", "auto"),
    ("generation.repair_passes", 500, 3),
    ("generation.outline_pass", "yes", True),
    ("generation.max_template_mb", 0, 32),
    ("generation.allow_user_uploads", "sure", True),
    ("generation.max_slides", 900, 50),
]:
    status, answered = call("PUT", "/admin/settings", {"values": {key: refused}}, expect=[422, 403])
    if status == 422:
        said = ((answered or {}).get("error") or {}).get("message") or ""
        if key not in said:
            failures.append(f"refusing {key} does not say which setting: {said!r}")
    call("PUT", "/admin/settings", {"values": {key: honoured}}, expect=[200, 403])

# Two slides headed the same thing. Everything the measurement looked at used to
# be about the drawing: a deck could be drawn perfectly with slide 5 and slide 8
# both headed "다음 단계", and nobody was told.
twice = data_of(call("POST", "/presentations", {"title": f"같은 제목 {RUN}", "prompt": "제목 반복 점검",
                                                "requestedSlideCount": 3, "language": "ko"}, expect=201))
call("PUT", f"/presentations/{twice['id']}/source", {"source": "\n".join([
    "# 다음 단계", "- 이관 일정을 확정합니다", "!notes 일정에 대해 말합니다",
    "# 비용", "- 올해 예산은 3억 원입니다", "!notes 예산에 대해 말합니다",
    "# 다음 단계", "- 담당자를 정합니다", "!notes 담당에 대해 말합니다"])}, expect=200)
measured = data_of(call("GET", f"/presentations/{twice['id']}/inspect", expect=200)) or {}
repeated = [f for f in (measured.get("findings") or []) if f.get("kind") == "twiceTitled"]
if not repeated:
    failures.append(f"two slides headed the same thing were not measured: {measured.get('findings')}")
elif repeated[0].get("slide") != 3 or "다음 단계" not in str(repeated[0].get("detail")):
    failures.append(f"the repeated heading was reported as {repeated[0]!r}")
call("DELETE", f"/presentations/{twice['id']}?permanent=true", expect=[204, 200])

# Rewriting a deck needs a connected model, and this deployment writes decks
# with the built-in writer. Asking for a rewrite anyway used to leave a deck
# somebody already had behind a failure screen, blaming an administrator for a
# service that is exactly as it ships.
polished = data_of(call("POST", "/presentations/generate",
                        {"prompt": f"다듬기 실패 확인 {RUN}", "slideCount": 3}, expect=[200, 201, 202])) or {}
polished_id = polished.get("id")
if polished_id:
    for _ in range(90):
        state = data_of(call("GET", f"/presentations/{polished_id}", expect=200)) or {}
        if state.get("status") in ("completed", "failed"):
            break
        time.sleep(2)
    if state.get("status") == "completed" and (state.get("slides") or []):
        # A deployment with no model host cannot rewrite an existing deck, and
        # says so instead of queueing work that will fail. Whichever way it
        # answers, the deck somebody already has must survive it.
        status, answered = call("POST", f"/presentations/{polished_id}/generate", {}, expect=[200, 202, 422])
        if status == 422:
            said = ((answered or {}).get("error") or {}).get("message") or ""
            if "모델" not in said:
                failures.append(f"the refusal does not say what a rewrite needs: {said!r}")
        else:
            for _ in range(90):
                after = data_of(call("GET", f"/presentations/{polished_id}", expect=200)) or {}
                if after.get("status") in ("completed", "failed"):
                    break
                time.sleep(2)
            if after.get("status") == "failed" and (after.get("slides") or []):
                failures.append("a deck with slides was left failed after a rewrite it could not do")
        kept = data_of(call("GET", f"/presentations/{polished_id}", expect=200)) or {}
        if not (kept.get("slides") or []):
            failures.append("a deck lost its slides to a rewrite that did not happen")
        # The polish button asks the same question at a different door, and used
        # to answer it by sending the author after an administrator.
        polish_status, polish_answer = call("POST", f"/presentations/{polished_id}/rewrite",
                                            {"instruction": "문장을 간결하게"}, expect=[200, 202, 409])
        if polish_status == 409:
            polished_error = ((polish_answer or {}).get("error") or {})
            if polished_error.get("code") not in ("rewrite_needs_model", "nothing_to_rewrite"):
                failures.append(f"the polish button refuses as {polished_error.get('code')!r}")
            if "관리자" in str(polished_error.get("message")):
                failures.append("the polish button sends the author after an administrator")
    call("DELETE", f"/presentations/{polished_id}?permanent=true", expect=[204, 200])

print("── templates ──")
templates = data_of(call("GET", "/templates?limit=100", expect=200)) or []
print(f"   {len(templates)} templates")
builtin = [t for t in templates if t["kind"] == "builtin"]
if len(builtin) < 50:
    failures.append(f"built-in template count is {len(builtin)}, expected 50")
# Another service picks a template before it asks for a deck, and a picker
# narrows before it scrolls. Both used to have to read the whole library.
narrowed = call("GET", "/templates?kind=builtin&limit=100", expect=200)
shipped = data_of(narrowed) or []
if not shipped or any(t["kind"] != "builtin" for t in shipped):
    failures.append("kind=builtin returned something that is not a built-in design")
own = data_of(call("GET", "/templates?kind=uploaded&limit=100", expect=200)) or []
if any(t["kind"] != "uploaded" for t in own):
    failures.append("kind=uploaded returned something that was not uploaded")
if len(shipped) + len(own) != len(templates):
    failures.append(f"the two kinds hold {len(shipped)}+{len(own)} of {len(templates)} templates")
named = data_of(call("GET", f"/templates?search={urllib.parse.quote(shipped[0]['name'])}", expect=200)) or []
if not any(t["id"] == shipped[0]["id"] for t in named):
    failures.append(f"searching for {shipped[0]['name']!r} did not find it")
if len(data_of(call("GET", "/templates?kind=nonsense&limit=100", expect=200)) or []) != len(templates):
    failures.append("an unknown kind narrowed the listing instead of being ignored")

# The whole point of listing them: name one, and the deck is built on it.
chosen = shipped[0]
on_template = data_of(call("POST", "/presentations/generate", {
    "title": f"템플릿 지정 {RUN}", "prompt": "사내 보안 교육 계획을 팀장들에게 보고합니다.",
    "language": "ko", "slideCount": 5, "templateId": chosen["id"]}, expect=202))
if (on_template or {}).get("templateId") != chosen["id"]:
    failures.append(f"a deck asked for on {chosen['name']} came back on {(on_template or {}).get('templateId')}")
else:
    for _ in range(120):
        built = data_of(call("GET", f"/presentations/{on_template['id']}", expect=200)) or {}
        if built.get("status") in ("ready", "completed", "failed"):
            break
        time.sleep(1)
    if built.get("templateName") != chosen["name"]:
        failures.append(f"the generated deck reports template {built.get('templateName')!r}")
    call("DELETE", f"/presentations/{on_template['id']}", expect=204)

first = templates[0]
call("GET", f"/templates/{first['id']}", expect=200)
call("GET", f"/templates/{first['id']}/preview.svg?width=320", raw=True, expect=200)
call("PUT", f"/templates/{first['id']}/favorite", {"favorite": True}, expect=200)
again = data_of(call("GET", "/templates?limit=100", expect=200)) or []
if not any(t["id"] == first["id"] and t.get("favorite") for t in again):
    failures.append("a starred template did not come back starred")
call("PUT", f"/templates/{first['id']}/favorite", {"favorite": False}, expect=200)
call("GET", "/templates/00000000-0000-0000-0000-000000000000", expect=404)
call("GET", "/templates/not-a-uuid", expect=400)

print("── presentations ──")
deck = data_of(call("POST", "/presentations", {
    "title": f"E2E 점검 {RUN}", "prompt": "전수 점검용 덱", "requestedSlideCount": 5, "language": "ko",
    "templateId": first["id"], "audience": "임원", "tone": "professional"}, expect=201))
deck_id = deck["id"]
call("GET", f"/presentations/{deck_id}", expect=200)
call("GET", "/presentations?limit=5", expect=200)
source = """# E2E 점검
@cover
> 전수 점검

# 지표
::kpi 핵심
- 처리량 | 1,200건
- 오류율 | 0.2%
::
- 첫 번째 요점
- 두 번째 요점

# 표
::table 분기 실적
- 항목 | 1분기 | 2분기
- 매출 | 120억 | 140억
::

# 마무리
@closing
> 감사합니다
"""
applied = data_of(call("PUT", f"/presentations/{deck_id}/source", {"source": source}, expect=200))
if not applied or not applied.get("applied"):
    failures.append("applying deck source did not report applied")
print("   applied:", (applied or {}).get("applied"), "findings:", len((applied or {}).get("findings") or []))
dry = data_of(call("PUT", f"/presentations/{deck_id}/source", {"source": source, "dryRun": True}, expect=200))
if dry and dry.get("applied"):
    failures.append("a dry run reported itself as applied")
call("GET", f"/presentations/{deck_id}/source", expect=200)
call("POST", f"/presentations/{deck_id}/source/preview.svg?slide=2&width=400", {"source": source}, raw=True, expect=200)
call("PUT", f"/presentations/{deck_id}/source", {"source": "# 한 장\n@cover\n", "version": 1}, expect=409, note="a stale version must conflict")
call("PUT", f"/presentations/{deck_id}/source", {"source": ""}, expect=422, note="empty source produces no slides")
call("PUT", f"/presentations/{deck_id}/source", {"source": "# x\n" * 40000}, expect=422, note="oversized source is refused")
# A model writes a directive against the next word when that word starts with a
# Korean particle. The reader has always taken !notes이 that way; every other
# directive has to be read the same, and a word that merely begins like one is
# not a directive at all.
glued = data_of(call("POST", "/presentations", {"title": f"붙여 쓴 지시어 {RUN}", "prompt": "점검"}, expect=201))
applied = data_of(call("PUT", f"/presentations/{glued['id']}/source", {"source":
    "# 제목\n@content\n- 요점\n!notes이 장은 핵심입니다\n!skip이 장은 부록입니다\n"}, expect=200)) or {}
if [w for w in (applied.get("warnings") or []) if "unknown directive" in w]:
    failures.append(f"a directive the reader knows was called unknown: {applied.get('warnings')}")
stored = data_of(call("GET", f"/presentations/{glued['id']}", expect=200)) or {}
glued_slide = (stored.get("slides") or [{}])[0]
if (glued_slide.get("speakerNotes") or "") != "이 장은 핵심입니다":
    failures.append(f"the note lost its first word: {glued_slide.get('speakerNotes')!r}")
if not (glued_slide.get("content") or {}).get("skipped"):
    failures.append("a slide told to be skipped was kept in the talk")
noted = data_of(call("PUT", f"/presentations/{glued['id']}/source", {"source":
    "# 제목\n@content\n- 요점\n!noteworthy 이건 지시어가 아닙니다\n"}, expect=200)) or {}
if not [w for w in (noted.get("warnings") or []) if "unknown directive" in w]:
    failures.append("a word that only begins like a directive was read as one")
call("DELETE", f"/presentations/{glued['id']}", expect=204)

# This deployment's cap on deck length is applied at the front door, and used to
# be applied silently: a request for more slides than are allowed came back 201
# with a shorter deck and no reason. The same cap on an imported file and on
# applied source both announce themselves.
print("── a deck shorter than what was asked for ──")
kept = {}
for row in (data_of(call("GET", "/admin/settings", expect=200)) or []):
    if row.get("key") in ("generation.max_slides", "generation.default_slide_count"):
        kept[row["key"]] = row.get("value")
call("PATCH", "/admin/settings", {"settings": {"generation.default_slide_count": 3}}, expect=200)
call("PATCH", "/admin/settings", {"settings": {"generation.max_slides": 5}}, expect=200)
capped = data_of(call("POST", "/presentations", {"title": "자름 점검", "prompt": "점검", "slideCount": 10}, expect=(201, 202)))
if capped:
    if (capped.get("requestedSlideCount") or 0) > 5:
        failures.append(f"the deployment allows 5 slides and made a deck of {capped.get('requestedSlideCount')}")
    told = data_of(call("GET", f"/presentations/{capped['id']}", expect=200)) or {}
    notes = told.get("generationNotes") or told.get("notes") or []
    if not any("10" in str(note) and "5" in str(note) for note in notes):
        failures.append(f"a deck cut from 10 slides to 5 says {notes}")
    inside = data_of(call("POST", "/presentations", {"title": "자름 점검 2", "prompt": "점검", "slideCount": 4}, expect=(201, 202)))
    if inside:
        quiet = data_of(call("GET", f"/presentations/{inside['id']}", expect=200)) or {}
        if quiet.get("generationNotes"):
            failures.append(f"a deck inside the cap was told {quiet['generationNotes']}")
        call("DELETE", f"/presentations/{inside['id']}", expect=(200, 204))
    call("DELETE", f"/presentations/{capped['id']}", expect=(200, 204))
for key in ("generation.max_slides", "generation.default_slide_count"):
    if kept.get(key) is not None:
        call("PATCH", "/admin/settings", {"settings": {key: kept[key]}}, expect=200)

# ── the exported file, read by somebody else ──
# Everything else here reads the export with our own code, which cannot tell us
# whether PowerPoint will make sense of it. python-pptx is an independent OOXML
# reader: what it can see is roughly what a person opening the file sees, and
# what it cannot see is not there for them either.
#
# The properties below are the ones that go silently wrong. A title drawn as a
# text box instead of a title placeholder leaves the outline view and the screen
# reader blank. A chart flattened to a picture cannot be edited or read aloud. A
# link written as blue text is not a link.
# ── every design this deployment ships ──
# A template is not one deck's problem: every deck made with it inherits
# whatever is wrong, and fifty of them travel inside the image. Nothing outside
# this product had ever opened them.
print("── every design this deployment ships ──")
try:
    from pptx import Presentation as DesignReader
except ImportError:
    print("   python-pptx is not installed here, so the designs were not opened independently")
    DesignReader = None
if DesignReader is not None:
    designs = data_of(call("GET", "/templates?limit=100", expect=200)) or []
    checks += 1
    if len(designs) < 1:
        failures.append("this deployment ships no designs at all")
    opened, sizes = 0, set()
    for design in designs:
        status, file = call("GET", f"/templates/{design['id']}/download", raw=True, expect=200)
        if status != 200 or not file:
            failures.append(f"the design {design['name']!r} could not be fetched")
            continue
        try:
            drawing = DesignReader(io.BytesIO(file))
        except Exception as error:
            failures.append(f"the design {design['name']!r} cannot be opened by an OOXML reader: {error}")
            continue
        opened += 1
        sizes.add((round(drawing.slide_width / 914400, 2), round(drawing.slide_height / 914400, 2)))
        layouts = list(drawing.slide_layouts)
        if design.get("layoutCount") and len(layouts) != design["layoutCount"]:
            failures.append(f"{design['name']!r}: the product says {design['layoutCount']} layouts "
                            f"and a reader sees {len(layouts)}")
        if drawing.slide_master is None:
            failures.append(f"{design['name']!r} has no slide master")
        roles = set()
        blank = 0
        for layout in layouts:
            if len(layout.placeholders) == 0:
                blank += 1
            for shape in layout.placeholders:
                roles.add(str(shape.placeholder_format.type))
        # The blank layout is meant to hold nothing; a second one is a place no
        # deck can be written into.
        if blank > 1:
            failures.append(f"{design['name']!r} has {blank} layouts with nowhere to put words")
        if not any(role.startswith("TITLE") or role.startswith("CENTER_TITLE") for role in roles):
            failures.append(f"{design['name']!r} offers no title placeholder on any layout")
        if not any(role.startswith(("BODY", "SUBTITLE", "OBJECT")) for role in roles):
            failures.append(f"{design['name']!r} offers nowhere for a deck's words")
    checks += 1
    if opened != len(designs):
        failures.append(f"{len(designs)} designs ship and {opened} could be opened")
    checks += 1
    if len(sizes) > 1:
        failures.append(f"the shipped designs are not one slide size: {sorted(sizes)}")

print("── the exported file, read by somebody else ──")
# The PDF gets the same treatment. The product reads its own PDFs back in its
# own tests, which cannot say whether anyone else's reader agrees: a file can
# render perfectly and still be unsearchable, unselectable and silent to a
# screen reader — drop the /ToUnicode map and it looks exactly the same on
# screen. pypdf is the independent judge.
try:
    from pypdf import PdfReader as ThirdPartyPDF
except ImportError:
    print("   pypdf is not installed here, so the PDF was not read back independently")
    ThirdPartyPDF = None
try:
    from pptx import Presentation as ThirdPartyReader
except ImportError:
    print("   python-pptx is not installed here, so the export was not read back independently")
    ThirdPartyReader = None
if ThirdPartyReader is not None:
    spec_deck = data_of(call("POST", "/presentations", {"title": f"규격 {RUN}", "prompt": "규격", "slideCount": 5}, expect=201)) or {}
    call("PUT", f"/presentations/{spec_deck['id']}/source", {"source":
        "# 표지가 되는 장\n@cover\n> 한 줄 리드\n\n"
        "# 분기 매출\n::columns 매출\n- 1분기 | 12\n- 2분기 | 15.5\n::\n\n"
        "# 분기별 비용\n::table 비용\n- 구분 | **1분기** | 비고\n- 인프라 | 1억 2천 | [계약서](https://example.invalid/contract)\n::\n\n"
        "# 링크가 있는 장\n- 자세한 내용은 [회사 안내](https://example.invalid/guide) 참고\n"
        "- 본체 크기는 1200*800*750 입니다\n"}, expect=200)
    _, exported = call("GET", f"/presentations/{spec_deck['id']}/export.pptx", raw=True, expect=200)
    try:
        show = ThirdPartyReader(io.BytesIO(exported or b""))
    except Exception as error:
        show = None
        failures.append(f"an independent OOXML reader could not open the exported file: {error}")
    if show is not None:
        checks += 1
        if len(show.slides) != 4:
            failures.append(f"the deck has 4 slides and the exported file reads as {len(show.slides)}")
        # shapes.title resolves from the placeholder index, so it answers even for
        # a title emitted as an ordinary body placeholder. The type is what the
        # outline view and a screen reader go by, so the type is what is checked.
        untitled = []
        for index, slide in enumerate(show.slides, start=1):
            kinds = {str(shape.placeholder_format.type) for shape in slide.shapes if shape.is_placeholder}
            if not any(kind.startswith("TITLE") or kind.startswith("CENTER_TITLE") for kind in kinds):
                untitled.append((index, sorted(kinds)))
        checks += 1
        if untitled:
            failures.append(f"these slides carry no title placeholder, so the outline view is blank for them: {untitled}")
        # "1200*800*750" is how a size is written; read as emphasis it drew
        # 1200800750, which is not markup on the wall but a different number.
        checks += 1
        drawn = " ".join(shape.text_frame.text for slide in show.slides
                         for shape in slide.shapes if shape.has_text_frame)
        if "1200*800*750" not in drawn:
            failures.append("a size written 1200*800*750 was not drawn as the author typed it")
        charts, tables, links = 0, 0, 0
        for slide in show.slides:
            for shape in slide.shapes:
                if shape.has_chart:
                    charts += 1
                    plot = shape.chart.plots[0]
                    if [str(one) for one in plot.categories] != ["1분기", "2분기"]:
                        failures.append(f"the exported chart's categories read as {[str(one) for one in plot.categories]}")
                    if [list(series.values) for series in plot.series] != [[12.0, 15.5]]:
                        failures.append("the exported chart's numbers are not the deck's")
                if getattr(shape, "has_table", False) and shape.has_table:
                    tables += 1
                    if [cell.text for cell in shape.table.rows[0].cells] != ["구분", "1분기", "비고"]:
                        failures.append("the exported table's first row is not the deck's")
                    # A cell was written as one piece of text, so a cell's marks
                    # were exported as the characters they are typed with: the
                    # file's header read **1분기** while the preview drew the
                    # word alone, and the address behind a cell was gone.
                    cells = [cell.text for row in shape.table.rows for cell in row.cells]
                    printed = [cell for cell in cells if "**" in cell or "](" in cell]
                    if printed:
                        failures.append(f"the exported table prints its markup where a reader sees it: {printed}")
                    linked = [run.hyperlink.address for row in shape.table.rows for cell in row.cells
                              for para in cell.text_frame.paragraphs for run in para.runs
                              if run.hyperlink is not None and run.hyperlink.address]
                    if linked != ["https://example.invalid/contract"]:
                        failures.append(f"the address written in a table cell reads as {linked}")
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.hyperlink is not None and run.hyperlink.address:
                                links += 1
        for what, found in (("chart", charts), ("table", tables), ("hyperlink", links)):
            checks += 1
            if found == 0:
                failures.append(f"the deck has a {what} and the exported file carries none a reader can find")
    if ThirdPartyPDF is not None:
        _, printed = call("GET", f"/presentations/{spec_deck['id']}/export.pdf", raw=True, expect=200)
        try:
            paper = ThirdPartyPDF(io.BytesIO(printed or b""))
        except Exception as error:
            paper = None
            failures.append(f"an independent PDF reader could not open the exported file: {error}")
        if paper is not None:
            checks += 1
            if len(paper.pages) != 4:
                failures.append(f"the deck has 4 slides and the PDF has {len(paper.pages)} pages")
            # A link written in a cell was live in the .pptx and dead in the PDF
            # of the same deck: the drawing carried the words alone.
            checks += 1
            clickable = []
            for page in paper.pages:
                for annotation in page.get("/Annots") or []:
                    entry = annotation.get_object()
                    action = entry.get("/A")
                    if entry.get("/Subtype") == "/Link" and action:
                        clickable.append(str(action.get("/URI")))
            if "https://example.invalid/contract" not in clickable:
                failures.append(f"a link written in a table cell is not clickable in the PDF: {clickable}")
            checks += 1
            wide = float(paper.pages[0].mediabox.width) / 72
            if not 13 < wide < 13.7:
                failures.append(f"the PDF's pages are {wide:.2f}in wide, not a 16:9 slide")
            lifted = "\n".join((page.extract_text() or "") for page in paper.pages)
            # Titles, a chart's own numbers, a table's cells and a link's words:
            # if any of these will not come out, the file cannot be searched.
            for wanted in ("표지가 되는 장", "분기 매출", "매출", "1분기", "인프라", "1억 2천", "회사 안내"):
                checks += 1
                if wanted not in lifted:
                    failures.append(f"a PDF reader cannot find {wanted!r} in the exported file")
    call("DELETE", f"/presentations/{spec_deck['id']}", expect=(200, 204))

print("── every run speaks the deck's own language ──")
# A screen reader picks its voice from the run's language and PowerPoint
# spell-checks against it; an English deck's table cells said ko-KR.
spoken_deck = data_of(call("POST", "/presentations", {"title": f"lang {RUN}", "prompt": "language", "slideCount": 3}, expect=201)) or {}
call("PUT", f"/presentations/{spoken_deck['id']}/source", {"source":
    "# Quarterly review\n- Throughput rose to 1,800 units an hour\n\n"
    "# Cost by quarter\n::table\n- Item | Q1\n- Infrastructure | 800\n::\n\n"
    "# Key numbers\n::kpi\n- Throughput | 1,800/h | +45%\n::\n\n"
    "# Trend\n::columns Throughput\n- Q1 | 1240\n::\n"}, expect=200)
call("PATCH", f"/presentations/{spoken_deck['id']}", {"language": "en"}, expect=[200, 204])
_, spoken_file = call("GET", f"/presentations/{spoken_deck['id']}/export.pptx", raw=True, expect=200)
with zipfile.ZipFile(io.BytesIO(spoken_file or b"")) as bundle:
    spoken_markup = "".join(bundle.read(name).decode("utf-8", "replace") for name in bundle.namelist()
                            if re.fullmatch(r"ppt/slides/slide\d+\.xml", name))
spoken_tags = sorted(set(re.findall(r'lang="([^"]+)"', spoken_markup)))
checks += 1
if spoken_tags != ["en-US"]:
    failures.append(f"an English deck declares its runs as {spoken_tags}")
call("DELETE", f"/presentations/{spoken_deck['id']}", expect=204)

print("── a picture keeps the words it already had ──")
# The alternative text somebody wrote in PowerPoint is the one thing about a
# picture nobody can work out by looking at it. Dropped on import, this product
# then asked the author to write what they had already written.
alt_pptx = None
try:
    from pptx import Presentation as AltWriter
    from pptx.util import Inches as AltInches
    written = AltWriter()
    AltWriter.__module__  # keep the import used
    cover = written.slides.add_slide(written.slide_layouts[5]); cover.shapes.title.text = "표지"
    shown = written.slides.add_slide(written.slide_layouts[5]); shown.shapes.title.text = "현장 사진"
    placed = shown.shapes.add_picture(io.BytesIO(png(rgb=(30, 140, 210))), AltInches(1), AltInches(2), AltInches(4), AltInches(3))
    placed._element._nvXxPr.cNvPr.set("descr", "자동 분류기가 상자를 컨베이어로 옮기는 모습")
    holder = io.BytesIO(); written.save(holder); alt_pptx = holder.getvalue()
except ImportError:
    print("   python-pptx is not installed here, so the picture's own words were not checked")
if alt_pptx:
    alt_deck = data_of(call("POST", "/presentations/import",
                            files={"file": (f"alt-{RUN}.pptx", alt_pptx, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}, expect=201)) or {}
    alt_id = (alt_deck.get("presentation") or alt_deck).get("id")
    alt_source = (data_of(call("GET", f"/presentations/{alt_id}/source", expect=200)) or {}).get("source", "")
    checks += 1
    if "자동 분류기가 상자를 컨베이어로 옮기는 모습" not in alt_source:
        failures.append("a picture's alternative text did not survive the import")
    checks += 1
    if [one for one in (data_of(call("GET", f"/presentations/{alt_id}/inspect", expect=200)) or {}).get("findings", [])
            if one.get("kind") == "undescribed"]:
        failures.append("an author was asked to describe a picture they had already described")
    call("DELETE", f"/presentations/{alt_id}", expect=204)
    # The import puts the picture in the library, where it stays after the deck
    # is gone; this run tidies up after itself.
    for one in data_of(call("GET", "/assets?limit=100", expect=200)) or []:
        if str(one.get("name", "")).startswith(f"alt-{RUN}"):
            call("DELETE", f"/assets/{one.get('id')}", expect=[204, 404])

print("── a picture with no words describing it ──")
# A picture is written into the file with its caption as alternative text and
# with nothing at all without one: a reader who cannot see the slide is told
# "Picture Placeholder 2", and nobody was asked for the words.
described_picture = data_of(call("POST", "/assets", files={"file": (f"alt-{RUN}.png", png(rgb=(120, 160, 90)), "image/png")},
                                 expect=201)) or {}
picture_name = described_picture.get("name")
if picture_name:
    described_deck = data_of(call("POST", "/presentations", {"title": f"대체텍스트 {RUN}", "prompt": "사진", "slideCount": 3}, expect=201)) or {}
    call("PUT", f"/presentations/{described_deck['id']}/source", {"source":
        "# 표지\n@cover\n> 한 줄\n\n"
        f"# 설명 없는 사진\n@content\n::image {picture_name}\n- 사진 옆의 글\n\n"
        f"# 설명 있는 사진\n@content\n::image {picture_name} | 자동 분류기가 상자를 옮기는 모습\n- 사진 옆의 글\n"}, expect=200)
    described_found = [one.get("slide") for one in (data_of(call("GET", f"/presentations/{described_deck['id']}/inspect", expect=200)) or {}).get("findings", [])
                       if one.get("kind") == "undescribed"]
    checks += 1
    if described_found != [2]:
        failures.append(f"a picture with no words is asked about on slides {described_found}, want only the one without a caption")
    _, described_file = call("GET", f"/presentations/{described_deck['id']}/export.pptx", raw=True, expect=200)
    checks += 1
    with zipfile.ZipFile(io.BytesIO(described_file or b"")) as bundle:
        drawn = "".join(bundle.read(name).decode("utf-8", "replace") for name in bundle.namelist()
                        if re.fullmatch(r"ppt/slides/slide\d+\.xml", name))
    if "자동 분류기가 상자를 옮기는 모습" not in drawn:
        failures.append("a picture's caption is not written into the file as its alternative text")
    call("DELETE", f"/presentations/{described_deck['id']}", expect=204)
    call("DELETE", f"/assets/{described_picture.get('id')}", expect=[204, 404])

print("── a jump the deck cannot follow is drawn as words ──")
# The exported file writes no relationship for it and the printed page no
# annotation; the drawing was the surface still promising something to click.
drawn_deck = data_of(call("POST", "/presentations", {"title": f"그린점프 {RUN}", "prompt": "점프", "slideCount": 3}, expect=201)) or {}
call("PUT", f"/presentations/{drawn_deck['id']}/source", {"source":
    "# 첫 장\n- 있는 장 [3장](#3) 과 없는 장 [9장 참고](#9)\n\n# 둘째 장\n- 2\n\n# 셋째 장\n- 3\n"}, expect=200)
_, drawn_svg = call("GET", f"/presentations/{drawn_deck['id']}/preview.svg?slide=1", raw=True, expect=200)
drawn_text = (drawn_svg or b"").decode("utf-8", "replace")
checks += 1
if 'href="#slide-9"' in drawn_text:
    failures.append("a jump past the end of the deck is drawn as something to click")
checks += 1
if 'href="#slide-3"' not in drawn_text or "9장 참고" not in re.sub(r"<[^>]+>", "", drawn_text):
    failures.append("a jump the deck can follow, or the words of one it cannot, went missing from the drawing")
call("DELETE", f"/presentations/{drawn_deck['id']}", expect=204)

print("── a jump on paper names the page it is printed on ──")
# The handout leaves out the slides being skipped, so the deck's third slide is
# not the third page; the jump was handed to the paper as the number alone.
paper_deck = data_of(call("POST", "/presentations", {"title": f"인쇄점프 {RUN}", "prompt": "인쇄", "slideCount": 3}, expect=201)) or {}
call("PUT", f"/presentations/{paper_deck['id']}/source", {"source":
    "# 첫 장\n- 세 번째 장 [3장 참고](#3) 과 없는 장 [9장](#9)\n\n# 둘째 장\n!skip\n- 뺀 장\n\n"
    "# 셋째 장\n- 여기가 3장\n\n# 넷째 장\n- 4장\n"}, expect=200)
_, printed_pdf = call("GET", f"/presentations/{paper_deck['id']}/export.pdf", raw=True, expect=200)
printed_body = (printed_pdf or b"").decode("latin-1")
checks += 1
landed = re.findall(r"/Dest \[(\d+) /Fit\]", printed_body)
if landed != ["1"]:
    failures.append(f"a jump on paper lands on pages {landed}, want the second page (index 1)")
checks += 1
if "/URI ()" in printed_body:
    failures.append("a jump past the end of the deck was printed as a link to nowhere")
call("DELETE", f"/presentations/{paper_deck['id']}", expect=204)

print("── a shared link says which slide of the deck each page is ──")
# A shared link leaves out the slides the author is skipping, so its third place
# is not the deck's third slide; a jump written [3장](#3) has to land on the one
# the deck numbers 3.
shared_deck = data_of(call("POST", "/presentations", {"title": f"공유점프 {RUN}", "prompt": "공유", "slideCount": 3}, expect=201)) or {}
call("PUT", f"/presentations/{shared_deck['id']}/source", {"source":
    "# 첫 장\n- 세 번째 장 [3장 참고](#3)\n\n# 둘째 장\n!skip\n- 뺀 장\n\n"
    "# 셋째 장\n- 여기가 3장\n\n# 넷째 장\n- 여기는 4장\n"}, expect=200)
shared_link = data_of(call("POST", f"/presentations/{shared_deck['id']}/shares", {"label": "검토", "days": 7}, expect=201)) or {}
shared_token = (shared_link.get("url") or "").rstrip("/").split("/")[-1]
shared_seen = data_of(call("GET", f"/shared/{shared_token}", expect=200, headers={})) or {}
checks += 1
if [one.get("position") for one in shared_seen.get("slides") or []] != [1, 3, 4]:
    failures.append(f"a shared link does not say which slide of the deck each page is: {shared_seen.get('slides')}")
call("DELETE", f"/presentations/{shared_deck['id']}", expect=204)

print("── a jump to a slide the deck does not have ──")
# A deck of two slides linking to its ninth wrote a relationship at
# ppt/slides/slide9.xml, a part that is not in the package it wrote.
jump_deck = data_of(call("POST", "/presentations", {"title": f"점프 {RUN}", "prompt": "점프", "slideCount": 2}, expect=201)) or {}
call("PUT", f"/presentations/{jump_deck['id']}/source", {"source":
    "# 첫 장\n- 뒤쪽 [9장 참고](#9) 과 있는 장 [2장](#2)\n\n# 둘째 장\n- 내용\n"}, expect=200)
_, jumped = call("GET", f"/presentations/{jump_deck['id']}/export.pptx", raw=True, expect=200)
with zipfile.ZipFile(io.BytesIO(jumped or b"")) as bundle:
    inside = set(bundle.namelist())
    dangling = []
    for entry in [name for name in inside if name.endswith(".rels")]:
        for target in re.findall(r'Type="[^"]*/slide"[^>]*Target="([^"]+)"', bundle.read(entry).decode("utf-8", "replace")):
            if "ppt/slides/" + target.split("/")[-1] not in inside:
                dangling.append((entry, target))
checks += 1
if dangling:
    failures.append(f"the exported package names slide parts it does not contain: {dangling}")
jump_found = [one for one in (data_of(call("GET", f"/presentations/{jump_deck['id']}/inspect", expect=200)) or {}).get("findings", [])
              if one.get("kind") == "link"]
checks += 1
if len(jump_found) != 1 or "9장 참고" not in jump_found[0].get("detail", ""):
    failures.append(f"a jump past the end of the deck is not reported: {jump_found}")
call("DELETE", f"/presentations/{jump_deck['id']}", expect=204)

print("── the front page counts without fetching everything ──")
# The front page worked its three numbers out by fetching every deck the account
# has, a hundred at a time.
summary_deck = data_of(call("POST", "/presentations", {"title": f"요약 {RUN}", "prompt": "요약", "slideCount": 1}, expect=201)) or {}
call("PUT", f"/presentations/{summary_deck['id']}/source", {"source": "# 한 장\n- 내용\n\n# 두 장\n- 내용\n"}, expect=200)
summary_now = data_of(call("GET", "/presentations/summary", expect=200)) or {}
summary_page = call("GET", "/presentations?limit=1&offset=0", expect=200)
summary_total = ((summary_page[1] or {}).get("meta") or {}).get("total") if isinstance(summary_page, tuple) else None
checks += 1
if not all(key in summary_now for key in ("total", "slides", "ready")):
    failures.append(f"the summary does not answer the three numbers the front page shows: {summary_now}")
checks += 1
if summary_total is not None and summary_now.get("total") != summary_total:
    failures.append(f"the summary counts {summary_now.get('total')} decks and the list says {summary_total}")
call("DELETE", f"/presentations/{summary_deck['id']}", expect=204)

print("── a run of guesses is written down once ──")
# The console recorded who signed in and nobody failing to. Writing down each
# attempt is not the answer: the limiter is in-process so that a failed sign-in
# is not a database write.
guess_before = [one for one in (data_of(call("GET", "/admin/audit?limit=100", expect=200)) or [])
                if one.get("action") == "auth.login_blocked"]
for guess in range(9):
    call("POST", "/auth/login", {"username": f"nobody-{RUN}@example.invalid", "password": f"wrong-{guess}"},
         expect=[401, 429], headers={})
guess_after = [one for one in (data_of(call("GET", "/admin/audit?limit=100", expect=200)) or [])
               if one.get("action") == "auth.login_blocked"]
checks += 1
if len(guess_after) - len(guess_before) != 1:
    failures.append(f"nine guesses wrote {len(guess_after) - len(guess_before)} rows, want exactly one")
checks += 1
if any("wrong-" in json.dumps(one) for one in guess_after):
    failures.append("the audit log holds what was typed as a password")

print("── opening and closing a link is written down ──")
# A share link hands the deck to whoever holds it, with no account and no
# sign-in. Creating an API key was written down and this was not.
audit_deck = data_of(call("POST", "/presentations", {"title": f"공유감사 {RUN}", "prompt": "감사", "slideCount": 1}, expect=201)) or {}
audit_share = data_of(call("POST", f"/presentations/{audit_deck['id']}/shares", {"label": f"검토 {RUN}", "days": 7}, expect=201)) or {}
audit_token = (audit_share.get("url") or "").rstrip("/").split("/")[-1]
call("DELETE", f"/presentations/{audit_deck['id']}/shares/{audit_share.get('id')}", expect=204)
audit_rows = data_of(call("GET", "/admin/audit?limit=50", expect=200)) or []
audit_rows = audit_rows if isinstance(audit_rows, list) else audit_rows.get("items", [])
audit_mine = [one for one in audit_rows if one.get("targetId") == audit_share.get("id")]
checks += 1
if sorted(one.get("action") for one in audit_mine) != ["share.create", "share.revoke"]:
    failures.append(f"opening and closing a link left {[one.get('action') for one in audit_mine]} in the audit log")
checks += 1
if audit_token and any(audit_token in json.dumps(one) for one in audit_rows):
    failures.append("the audit log holds the token that opens a shared deck")
call("DELETE", f"/presentations/{audit_deck['id']}", expect=204)

print("── both bar charts, each the way round it is named ──")
# The guide named the horizontal bar chart under none of the four names that
# reach it, and used `bars` — a spelling of the vertical one — for it.
bar_deck = data_of(call("POST", "/presentations", {"title": f"막대 {RUN}", "prompt": "막대", "slideCount": 2}, expect=201)) or {}
call("PUT", f"/presentations/{bar_deck['id']}/source", {"source":
    "# 세로\n::columns 매출\n- 서울 | 1240\n- 부산 | 620\n::\n\n"
    "# 가로\n::hbars 매출\n- 서울 | 1240\n- 부산 | 620\n::\n"}, expect=200)
for bar_slide, bar_axis in ((1, "wide"), (2, "tall")):
    _, bar_svg = call("GET", f"/presentations/{bar_deck['id']}/preview.svg?slide={bar_slide}", raw=True, expect=200)
    drawn = (bar_svg or b"").decode("utf-8", "replace")
    # the baseline: a column chart rules along the foot, a bar chart down the side
    rules = [(float(w), float(h)) for w, h in
             re.findall(r'<rect[^>]*width="([\d.]+)"[^>]*height="([\d.]+)"[^>]*rx="0\.0"', drawn)]
    checks += 1
    if not rules:
        failures.append(f"slide {bar_slide} drew no chart axis at all")
    elif (rules[0][0] > rules[0][1]) != (bar_axis == "wide"):
        failures.append(f"slide {bar_slide} drew its bars the wrong way round: axis {rules[0]}")
call("DELETE", f"/presentations/{bar_deck['id']}", expect=204)

print("── what is cut is said, and what fits is drawn whole ──")
# Everything a component can drop has to be reported, and everything it can
# hold has to arrive: a cell that wraps inside its row is not trimmed, and an
# address with nowhere to break is not cut in half.
cut_cases = {
    "table rows": "# 비용\n::table\n- 구분 | 값\n" + "".join(f"- 항목{n} | {n}00\n" for n in range(1, 13)) + "::\n",
    "table columns": "# 비용\n::table\n- 구분 | 1분기 | 2분기 | 3분기 | 4분기 | 상반기 | 비고\n- 설비 | 1 | 2 | 3 | 4 | 5 | 계약\n- 인건비 | 1 | 2 | 3 | 4 | 5 | 유지\n::\n",
    "steps entries": "# 순서\n::steps\n" + "".join(f"- 단계{n} | {n}월까지\n" for n in range(1, 10)) + "::\n",
}
for what, cut_source in cut_cases.items():
    cut_one = data_of(call("POST", "/presentations", {"title": f"잘림 {RUN}", "prompt": "잘림", "slideCount": 1}, expect=201)) or {}
    call("PUT", f"/presentations/{cut_one['id']}/source", {"source": cut_source}, expect=200)
    checks += 1
    if not [one for one in (data_of(call("GET", f"/presentations/{cut_one['id']}/inspect", expect=200)) or {}).get("findings", [])
            if one.get("kind") == "trimmed"]:
        failures.append(f"a component dropped {what} and said nothing")
    call("DELETE", f"/presentations/{cut_one['id']}", expect=204)

# And the other way: a long cell and an address that cannot be broken both
# arrive whole, so nothing is reported and nothing is lost.
long_cell = "아주 긴 설명을 한 칸에 넣습니다 " * 10
long_link = "https://intranet.example.invalid/reports/2026/logistics/automation/quarterly.pdf"
whole_one = data_of(call("POST", "/presentations", {"title": f"온전 {RUN}", "prompt": "온전", "slideCount": 1}, expect=201)) or {}
call("PUT", f"/presentations/{whole_one['id']}/source", {"source":
    f"# 비용\n::table\n- 구분 | 설명\n- 설비 | {long_cell}\n- 링크 | {long_link}\n::\n"}, expect=200)
_, whole_file = call("GET", f"/presentations/{whole_one['id']}/export.pptx", raw=True, expect=200)
with zipfile.ZipFile(io.BytesIO(whole_file or b"")) as bundle:
    whole_markup = "".join(bundle.read(name).decode("utf-8", "replace") for name in bundle.namelist()
                           if re.fullmatch(r"ppt/slides/slide\d+\.xml", name))
checks += 1
if long_link not in whole_markup:
    failures.append("an address with nowhere to break did not reach the slide whole")
checks += 1
if long_cell.strip() not in whole_markup:
    failures.append("a long cell that fits its row did not reach the slide whole")
checks += 1
if [one for one in (data_of(call("GET", f"/presentations/{whole_one['id']}/inspect", expect=200)) or {}).get("findings", [])
        if one.get("kind") == "trimmed"]:
    failures.append("nothing was dropped and the deck was told something had been")
call("DELETE", f"/presentations/{whole_one['id']}", expect=204)

print("── a heading cut on a connective ──")
# "비용을 줄이고" is one clause of two. 보고, 참고, 사고, 창고 and 광고 all end a
# heading perfectly well and every one of them ends in 고.
cut_deck = data_of(call("POST", "/presentations", {"title": f"제목 {RUN}", "prompt": "제목", "slideCount": 2}, expect=201)) or {}
call("PUT", f"/presentations/{cut_deck['id']}/source", {"source":
    "# 비용을 줄이고\n- 내용 한 줄\n\n# 분기 실적 보고\n- 내용 한 줄\n"}, expect=200)
cut_found = [one.get("slide") for one in (data_of(call("GET", f"/presentations/{cut_deck['id']}/inspect", expect=200)) or {}).get("findings", [])
             if one.get("kind") == "unfinished"]
checks += 1
if cut_found != [1]:
    failures.append(f"a heading cut on a connective is reported on slides {cut_found}, want only the cut one")
call("DELETE", f"/presentations/{cut_deck['id']}", expect=204)

print("── the same point made twice, in two forms of the same words ──")
# Korean puts an ending on both sides far more often than it leaves one bare,
# and a word was only ever matched against its own bare stem.
twice_deck = data_of(call("POST", "/presentations", {"title": f"반복 {RUN}", "prompt": "반복", "slideCount": 2}, expect=201)) or {}
call("PUT", f"/presentations/{twice_deck['id']}/source", {"source":
    "# 도입 효과\n- 자동화를 도입하면 오출고율이 크게 낮아집니다\n"
    "- 오출고율은 자동화 도입으로 크게 낮아집니다\n\n"
    "# 분기 실적\n- 매출은 전년 대비 12% 늘었습니다\n- 비용은 전년 대비 8% 줄었습니다\n"}, expect=200)
twice_found = [one.get("slide") for one in (data_of(call("GET", f"/presentations/{twice_deck['id']}/inspect", expect=200)) or {}).get("findings", [])
               if one.get("kind") == "repeat"]
checks += 1
if twice_found != [1]:
    failures.append(f"the same point written twice is reported on slides {twice_found}, want only the slide that says it twice")
call("DELETE", f"/presentations/{twice_deck['id']}", expect=204)

print("── a roadmap written as steps is read as a plan ──")
# A step and its date are one line to anybody looking at the slide, and were
# walked as the two strings they are stored as.
plan_deck = data_of(call("POST", "/presentations", {"title": f"계획 {RUN}", "prompt": "계획", "slideCount": 2}, expect=201)) or {}
call("PUT", f"/presentations/{plan_deck['id']}/source", {"source":
    "# 이행 계획\n::steps\n- 설계 | 2019년 3월까지\n- 설치 | 2019년 6월까지\n::\n\n"
    "# 끝난 일\n::steps\n- 설계 | 2019년 3월에 마쳤습니다\n- 설치 | 2019년 6월에 끝냈습니다\n::\n"}, expect=200)
plan_found = [one.get("slide") for one in (data_of(call("GET", f"/presentations/{plan_deck['id']}/inspect", expect=200)) or {}).get("findings", [])
              if one.get("kind") == "stale"]
checks += 1
if plan_found != [1]:
    failures.append(f"a roadmap dated to a year long past is reported on slides {plan_found}, want only the one that is still a plan")
call("DELETE", f"/presentations/{plan_deck['id']}", expect=204)

print("── a chart is asked where its numbers came from ──")
# The same figures asked about in a point, a table and a KPI row went unasked
# as a bar chart, which is the thing a room quotes.
chart_deck = data_of(call("POST", "/presentations", {"title": f"차트 {RUN}", "prompt": "차트", "slideCount": 2}, expect=201)) or {}
call("PUT", f"/presentations/{chart_deck['id']}/source", {"source":
    "# 처리량 추이\n::columns 처리량\n- 1분기 | 1240\n- 2분기 | 1520\n::\n\n"
    "# 출처를 밝힌 장\n::columns 처리량\n- 1분기 | 1240\n::\n!source 물류본부 2026 집계\n"}, expect=200)
chart_found = [one for one in (data_of(call("GET", f"/presentations/{chart_deck['id']}/inspect", expect=200)) or {}).get("findings", [])
               if one.get("kind") == "source"]
checks += 1
if len(chart_found) != 1 or "1240" not in chart_found[0].get("detail", ""):
    failures.append(f"a chart's numbers are not asked where they came from: {chart_found}")
checks += 1
if any(one.get("slide") == 2 for one in chart_found):
    failures.append("a chart that cites its source is still asked for one")
call("DELETE", f"/presentations/{chart_deck['id']}", expect=204)

print("── a refused link, wherever it is written ──")
# The same mistake in the same deck got two answers: reported in a point, and
# silent in a cell of a table, where it draws its markup on the wall all the same.
refused_deck = data_of(call("POST", "/presentations", {"title": f"링크 {RUN}", "prompt": "링크", "slideCount": 3}, expect=201)) or {}
call("PUT", f"/presentations/{refused_deck['id']}/source", {"source":
    "# 줄에 쓴 것\n- 자세한 내용은 [회사 안내](www.example.com) 참고\n\n"
    "# 표에 쓴 것\n::table\n- 항목 | 근거\n- 인프라 | [계약서](www.example.invalid/c)\n::\n"}, expect=200)
refused_found = [one for one in (data_of(call("GET", f"/presentations/{refused_deck['id']}/inspect", expect=200)) or {}).get("findings", [])
                 if one.get("kind") == "link"]
checks += 1
if len(refused_found) != 2:
    failures.append(f"a refused link is reported in a point and not in a table cell: {refused_found}")
call("DELETE", f"/presentations/{refused_deck['id']}", expect=204)

print("── inspect, preview, export ──")

# The path may name the format and the query may override it. The default used
# to be written between the two, so the branch that read .pdf could never run —
# /export.pdf was not even routed, and would have handed back a pptx if it had
# been. Both suffixes name what they say now.
for suffix, want in [("/export.pptx", b"PK"), ("/export.pdf", b"%PDF"), ("/export", b"PK")]:
    status, body = call("GET", f"/presentations/{deck_id}{suffix}", raw=True, expect=200)
    if status == 200 and not (body or b"").startswith(want):
        failures.append(f"{suffix} returned {(body or b'')[:4]!r}, expected a file starting {want!r}")
status, _ = call("GET", f"/presentations/{deck_id}/export?format=svg", expect=422,
                 note="an export format the product does not have")

inspected = data_of(call("GET", f"/presentations/{deck_id}/inspect", expect=200))
print("   slides:", (inspected or {}).get("slides"), "defects:", (inspected or {}).get("defects"),
      "advisories:", (inspected or {}).get("advisories"))
if (inspected or {}).get("defects"):
    failures.append(f"a plain four-slide deck has defects: {inspected}")
# A component draws as many entries as its kind can show and takes the first of
# them. A slide holding more than that has entries on no slide at all, and the
# measurement has to say so — a deck once asked a board to approve a plan whose
# "calculate the financial impact" step was drawn nowhere.
trimmed_deck = data_of(call("POST", "/presentations", {"title": f"안 그려지는 항목 {RUN}", "prompt": "단계 여섯"}, expect=201))
call("PUT", f"/presentations/{trimmed_deck['id']}/source", {"source":
     "# 이행 순서\n@content\n::steps 순서\n- 1 | 현황 파악\n- 2 | 대안 정의\n- 3 | 비교\n"
     "- 4 | 이행 계획\n- 5 | 재무 영향\n- 6 | 승인 요청\n::\n"}, expect=200)
measured = data_of(call("GET", f"/presentations/{trimmed_deck['id']}/inspect", expect=200)) or {}
carried = [f for f in (measured.get("findings") or []) if f.get("kind") == "trimmed"]
if not carried:
    failures.append("a six-entry component drawn five at a time is not reported")
elif "5 of its 6" not in carried[0].get("detail", ""):
    failures.append(f"the finding does not say what was left out: {carried[0]}")
# Split as the workspace's safe fix splits it, and the deck measures clean.
call("PUT", f"/presentations/{trimmed_deck['id']}/source", {"source":
     "# 이행 순서\n@content\n::steps 순서\n- 1 | 현황 파악\n- 2 | 대안 정의\n- 3 | 비교\n"
     "- 4 | 이행 계획\n- 5 | 재무 영향\n::\n\n# 이행 순서 (계속)\n@content\n::steps 순서\n- 6 | 승인 요청\n::\n"}, expect=200)
after = data_of(call("GET", f"/presentations/{trimmed_deck['id']}/inspect", expect=200)) or {}
if [f for f in (after.get("findings") or []) if f.get("kind") == "trimmed"]:
    failures.append("carrying the rest onto a second slide did not settle the finding")
call("DELETE", f"/presentations/{trimmed_deck['id']}", expect=204)

# A slide kept for the questions afterwards is walked past in the talk and is
# still in the deck and in the file. The flag rides in the slide's content, so
# what this is really checking is that it survives the four hands it passes
# through: the source language, the stored slide, the file, and the source
# written back out.
skip_deck = data_of(call("POST", "/presentations", {"title": f"건너뛰는 장 {RUN}", "prompt": "부록"}, expect=201))
call("PUT", f"/presentations/{skip_deck['id']}/source", {"source":
     "# 요약\n@cover\n> 올해 계획\n\n# 부록: 산출 근거\n@content\n!skip\n- 계산 방법\n"}, expect=200)
kept = data_of(call("GET", f"/presentations/{skip_deck['id']}", expect=200)) or {}
marks = [bool((slide.get("content") or {}).get("skipped")) for slide in kept.get("slides") or []]
if marks != [False, True]:
    failures.append(f"the skipped slide is not marked in the stored deck: {marks}")
written = (data_of(call("GET", f"/presentations/{skip_deck['id']}/source", expect=200)) or {}).get("source", "")
if "!skip" not in written:
    failures.append("the source written back out lost the skip")
status, skip_pptx = call("GET", f"/presentations/{skip_deck['id']}/export?format=pptx", raw=True, expect=200)
try:
    package = zipfile.ZipFile(io.BytesIO(skip_pptx))
    hidden = [name for name in package.namelist() if name.startswith("ppt/slides/slide")
              and b'show="0"' in package.read(name)]
    if len(hidden) != 1:
        failures.append(f"the exported file marks {len(hidden)} slides hidden, expected 1")
except Exception as error:
    failures.append(f"the exported pptx could not be read: {error}")
call("DELETE", f"/presentations/{skip_deck['id']}", expect=204)

# A link is written in the text and drawn as words. What the sweep is watching
# for is the markup reaching a screen: the preview and the file both have to
# draw 안내 문서, and neither may draw the brackets or the address.
link_deck = data_of(call("POST", "/presentations", {"title": f"링크 {RUN}", "prompt": "안내"}, expect=201))
call("PUT", f"/presentations/{link_deck['id']}/source", {"source":
     "# 안내\n@content\n- 자세한 내용은 [안내 문서](https://docs.example.com/plan)를 보십시오\n"
     "- 근거는 [부록](#2)에 있습니다\n\n# 부록\n@content\n- 산출 근거\n"}, expect=200)
status, preview = call("GET", f"/presentations/{link_deck['id']}/preview.svg?slide=1&width=900", raw=True, expect=200)
drawn = preview.decode("utf-8", "replace")
# The address belongs in the link the drawing carries, never in the words it
# draws: what is drawn is what sits between the tags.
if ">https://docs.example.com/plan" in drawn or "](" in drawn:
    failures.append("the preview draws the link markup")
if '<a href="https://docs.example.com/plan"' not in drawn:
    failures.append("the drawing does not carry the link where presenting can follow it")
if "안내 문서" not in drawn or "text-decoration=\"underline\"" not in drawn:
    failures.append("the preview does not draw the link as one")
status, linked_pptx = call("GET", f"/presentations/{link_deck['id']}/export?format=pptx", raw=True, expect=200)
try:
    package = zipfile.ZipFile(io.BytesIO(linked_pptx))
    slide = package.read("ppt/slides/slide1.xml").decode("utf-8")
    rels = package.read("ppt/slides/_rels/slide1.xml.rels").decode("utf-8")
    if "hlinkClick" not in slide or "](" in slide:
        failures.append("the exported slide does not carry the link")
    if "https://docs.example.com/plan" not in rels or 'TargetMode="External"' not in rels:
        failures.append("the link is not a relationship of the exported slide")
    if "slide2.xml" not in rels:
        failures.append("the jump to another slide is not a relationship of the exported slide")
except Exception as error:
    failures.append(f"the exported pptx could not be read: {error}")
# Presenting draws the slide itself so a link on it can be clicked; a jump has
# to name the slide it goes to.
if '<a href="#slide-2"' not in drawn:
    failures.append("the drawing does not name the slide a jump goes to")
# A target the deck will not follow draws as its own markup, and is reported.
call("PUT", f"/presentations/{link_deck['id']}/source", {"source":
     "# 안내\n@content\n- [문서](www.example.com)를 보십시오\n"}, expect=200)
measured = data_of(call("GET", f"/presentations/{link_deck['id']}/inspect", expect=200)) or {}
if not [f for f in (measured.get("findings") or []) if f.get("kind") == "link"]:
    failures.append("a link the deck cannot follow is not reported")
# A link written in the speaker notes has to keep its address: the notes are
# where the presenter wrote it down, and a link that draws its words and loses
# where it points is worse than no link at all.
notes_deck = data_of(call("POST", "/presentations", {"title": f"노트 링크 {RUN}", "prompt": "안내"}, expect=201))
call("PUT", f"/presentations/{notes_deck['id']}/source", {"source":
     "# 안내\n@content\n- 요점\n!notes 근거는 [분기 보고서](https://reports.example.com/q3)에 있습니다\n"}, expect=200)
status, notes_pptx = call("GET", f"/presentations/{notes_deck['id']}/export?format=pptx", raw=True, expect=200)
try:
    package = zipfile.ZipFile(io.BytesIO(notes_pptx))
    notes = package.read("ppt/notesSlides/notesSlide1.xml").decode("utf-8")
    notes_rels = package.read("ppt/notesSlides/_rels/notesSlide1.xml.rels").decode("utf-8")
    if "분기 보고서" not in notes or "](" in notes:
        failures.append("the notes do not draw the words of the link")
    if "https://reports.example.com/q3" not in notes_rels:
        failures.append("the notes lost the address the link points at")
except Exception as error:
    failures.append(f"the exported notes could not be read: {error}")
call("DELETE", f"/presentations/{notes_deck['id']}", expect=204)

call("DELETE", f"/presentations/{link_deck['id']}", expect=204)

# A word marked inside a line: the file carries the weight, the preview shows
# it, and neither draws the stars.
mark_deck = data_of(call("POST", "/presentations", {"title": f"강조 {RUN}", "prompt": "실적"}, expect=201))
call("PUT", f"/presentations/{mark_deck['id']}/source", {"source":
     "# 실적\n@content\n- 이번 분기 **매출이 12% 늘었습니다**\n- *가정은 인건비 동결입니다*\n- 3 * 4 = 12\n"}, expect=200)
status, marked_svg = call("GET", f"/presentations/{mark_deck['id']}/preview.svg?slide=1&width=900", raw=True, expect=200)
drawn = marked_svg.decode("utf-8", "replace")
if "**" in drawn or "매출이 12% 늘었습니다" not in drawn:
    failures.append("the preview draws the marks instead of the words")
if 'font-weight="700"' not in drawn or 'font-style="italic"' not in drawn:
    failures.append("the preview does not show the marked words as marked")
if "3 * 4 = 12" not in drawn:
    failures.append("a star with no partner was eaten")
status, marked_pptx = call("GET", f"/presentations/{mark_deck['id']}/export?format=pptx", raw=True, expect=200)
try:
    slide = zipfile.ZipFile(io.BytesIO(marked_pptx)).read("ppt/slides/slide1.xml").decode("utf-8")
    if "**" in slide or 'b="1"' not in slide or 'i="1"' not in slide:
        failures.append("the exported slide does not carry the marks")
except Exception as error:
    failures.append(f"the exported pptx could not be read: {error}")
call("DELETE", f"/presentations/{mark_deck['id']}", expect=204)

call("GET", f"/presentations/{deck_id}/preview.svg?slide=1&width=400", raw=True, expect=200)
status, pptx = call("GET", f"/presentations/{deck_id}/export?format=pptx", raw=True, expect=200)
try:
    package = zipfile.ZipFile(io.BytesIO(pptx))
    slides = [n for n in package.namelist() if n.startswith("ppt/slides/slide")]
    print(f"   pptx {len(pptx):,} bytes, {len(slides)} slides")
    if len(slides) != 4:
        failures.append(f"exported {len(slides)} slides, expected 4")
    bad = package.testzip()
    if bad:
        failures.append(f"the exported package is corrupt at {bad}")
except Exception as error:
    failures.append(f"the exported pptx could not be read: {error}")
# The deck on paper: one page per slide, at the deck's own size, with the words
# where the screen puts them and the links still live.
# Which of this deployment's templates paints a wash, if any.
wash_template = next((t.get("id") for t in (data_of(call("GET", "/templates?limit=50", expect=200)) or [])
                      if "Wash" in (t.get("name") or "")), "")
# The built-in PDF face covers Korean, Latin and kana. A deck written in
# Japanese 新字体 or simplified Chinese reaches past it and those characters
# print as blank space, so the export has to say which ones it dropped: finding
# out here is the difference from finding out in print.
# A heading longer than its box is cut with an ellipsis, and the words past the
# cut are on no slide. Losing them is sometimes the only thing to do; not saying
# so is not: sixteen words were written, eight were drawn, and the quality panel
# called the deck clean.
print("── a heading that had to be cut says so ──")
long_heading = " ".join(f"항목{n:02d}가나" for n in range(1, 17))
cut = data_of(call("POST", "/presentations", {"title": "잘림 점검", "prompt": "점검", "language": "ko"}, expect=201))
if cut:
    applied = data_of(call("PUT", f"/presentations/{cut['id']}/source",
                           {"source": f"# {long_heading}\n- 짧은 요점\n"}, expect=200)) or {}
    said = " | ".join(applied.get("warnings") or [])
    checks += 1
    if "shortened to fit" not in said:
        failures.append(f"a heading was cut and the apply said nothing: {said!r}")
    measured = data_of(call("GET", f"/presentations/{cut['id']}/inspect", expect=200)) or {}
    kinds = {f.get("kind") for f in (measured.get("findings") or [])}
    checks += 1
    if "trimmed" not in kinds:
        failures.append(f"a heading was cut and the inspector did not report it: {sorted(kinds)}")
    print(f"   warned: {'shortened to fit' in said} · reported: {'trimmed' in kinds}")
    # And a heading that fits stays quiet on both.
    whole = data_of(call("PUT", f"/presentations/{cut['id']}/source",
                         {"source": "# 물류센터 자동화 도입 승인\n- 짧은 요점\n"}, expect=200)) or {}
    quiet = data_of(call("GET", f"/presentations/{cut['id']}/inspect", expect=200)) or {}
    checks += 1
    if "trimmed" in {f.get("kind") for f in (quiet.get("findings") or [])}:
        failures.append("a heading that fits was reported as cut")
    checks += 1
    if "shortened to fit" in " | ".join(whole.get("warnings") or []):
        failures.append("a heading that fits was warned about")
    call("DELETE", f"/presentations/{cut['id']}", expect=204)

print("── what the pdf face cannot draw ──")
japanese = data_of(call("POST", "/presentations", {"title": "四半期実績", "prompt": "四半期実績の報告。",
                                                   "language": "ja", "requestedSlideCount": 3}, expect=201))
if japanese:
    call("PUT", f"/presentations/{japanese['id']}/source",
         {"source": "# 四半期実績\n- 売上と直販の実績\n- 业绩と销售额\n"}, expect=200)


    def export_headers(deck):
        request = urllib.request.Request(f"{BASE}/presentations/{deck}/export?format=pdf",
                                         headers={"X-Ptium-Dev-Secret": SECRET})
        with urllib.request.urlopen(request) as response:
            return response.headers.get("X-Ptium-Missing-Characters") or "", response.read()

    reported, printed = export_headers(japanese["id"])
    reported = urllib.parse.unquote_plus(reported)
    checks += 1
    if not printed.startswith(b"%PDF-"):
        failures.append("a Japanese deck did not export as a PDF at all")
    for character in "実売业":
        checks += 1
        if character not in reported:
            failures.append(f"{character!r} prints as a blank and the export did not say so: {reported!r}")
    print(f"   reported as undrawable: {reported!r}")
    korean, _ = export_headers(deck_id)
    checks += 1
    if korean:
        failures.append(f"a Korean deck reported characters it could draw: {korean!r}")
    call("DELETE", f"/presentations/{japanese['id']}", expect=204)

status, pdf_bytes = call("GET", f"/presentations/{deck_id}/export?format=pdf", raw=True, expect=200)
if not pdf_bytes.startswith(b"%PDF-") or not pdf_bytes.rstrip().endswith(b"%%EOF"):
    failures.append("the PDF export is not a PDF")
else:
    body = pdf_bytes.decode("latin-1")
    pages = body.count("/Type /Page ")
    if pages != 4:
        failures.append(f"a four-slide deck printed {pages} pages")
    if "/Subtype /Type0" not in body or "/FontFile2" not in body:
        failures.append("the PDF does not carry the font its text is set in")
    if "/ToUnicode" not in body:
        failures.append("the PDF's text cannot be copied out as words")
    # A template whose background is a wash prints as one: resolving it to a
    # flat colour would be a page that disagrees with every other drawing of the
    # same slide.
    wash = data_of(call("POST", "/presentations", {"title": f"워시 {RUN}", "prompt": "점검",
                                                   "templateId": wash_template}, expect=201)) if wash_template else None
    if wash:
        call("PUT", f"/presentations/{wash['id']}/source", {"source": "# 표지입니다\n@cover\n> 부제목\n"}, expect=200)
        status, washed = call("GET", f"/presentations/{wash['id']}/export?format=pdf", raw=True, expect=200)
        if b"/ShadingType 2" not in washed:
            failures.append("a wash on the screen printed as a flat colour")
        call("DELETE", f"/presentations/{wash['id']}", expect=204)
    print(f"   pdf {len(pdf_bytes):,} bytes, {pages} pages")
    # The same picture on several slides is one picture in the file. A deck's
    # logo carried on every slide made a forty-page PDF nine times heavier.
    logo = data_of(call("POST", "/assets", files={"file": (f"repeat-{RUN}.png", png(rgb=(40, 90, 200)), "image/png")},
                        expect=201)) or {}
    picture_deck = data_of(call("POST", "/presentations", {"title": f"같은 그림 {RUN}", "prompt": "점검"}, expect=201))
    call("PUT", f"/presentations/{picture_deck['id']}/source", {"source":
         "".join(f"# {n}장\n@content\n- 요점\n::image {logo.get('name')} | 현장\n\n" for n in range(1, 6))}, expect=200)
    status, repeated = call("GET", f"/presentations/{picture_deck['id']}/export?format=pdf", raw=True, expect=200)
    stored = repeated.count(b"/Subtype /Image")
    if stored != 1:
        failures.append(f"one picture on five slides is stored {stored} times")
    call("DELETE", f"/presentations/{picture_deck['id']}", expect=204)
    call("DELETE", f"/assets/{logo.get('id')}", expect=204)
# A deck's own source, applied back to it, changes nothing.
#
# It changed two of the twelve real decks measured against this product: a slide
# holding a picture beside its prose exchanged the two on every save, because the
# source was written in the order the slide reads and taken back in the order the
# compiler fills regions. Nothing here noticed — every check asked what a deck
# holds, and none asked whether writing it down and reading it back gives the
# same deck.
print("── a deck's own source, applied back ──")
left = data_of(call("POST", "/assets", files={"file": (f"rt-left-{RUN}.png", png(rgb=(200, 60, 60)), "image/png")},
                    expect=201)) or {}
right = data_of(call("POST", "/assets", files={"file": (f"rt-right-{RUN}.png", png(rgb=(60, 60, 200)), "image/png")},
                     expect=201)) or {}
round_trip = data_of(call("POST", "/presentations", {"title": f"왕복 {RUN}", "prompt": "점검"}, expect=201))
call("PUT", f"/presentations/{round_trip['id']}/source", {"source":
     f"# 두 그림이 있는 장\n@content\n- 왼쪽에 붙는 그림과 오른쪽에 붙는 그림\n"
     f"::image {left.get('name')} | 왼쪽\n::image {right.get('name')} | 오른쪽\n"
     f"\n# 다음 장\n- 한 줄\n"}, expect=200)
before_deck = data_of(call("GET", f"/presentations/{round_trip['id']}")) or {}
written = data_of(call("GET", f"/presentations/{round_trip['id']}/source")) or {}
source_text = written.get("source") if isinstance(written, dict) else written
checks += 1
if not str(source_text or "").strip():
    failures.append("a deck's source came back empty")
# Twice, because the order a slide's pictures are written in was decided by
# walking a map: the same deck gave two different documents on two reads.
second = data_of(call("GET", f"/presentations/{round_trip['id']}/source")) or {}
checks += 1
if (second.get("source") if isinstance(second, dict) else second) != source_text:
    failures.append("reading the same deck's source twice gave two different documents")
call("PUT", f"/presentations/{round_trip['id']}/source", {"source": source_text}, expect=200)
after_deck = data_of(call("GET", f"/presentations/{round_trip['id']}")) or {}
def slide_contents(deck):
    return [json.dumps(slide.get("content"), ensure_ascii=False, sort_keys=True)
            for slide in (deck.get("slides") or [])]
checks += 1
if slide_contents(before_deck) != slide_contents(after_deck):
    moved = [n + 1 for n, (b, a) in enumerate(zip(slide_contents(before_deck), slide_contents(after_deck))) if b != a]
    failures.append(f"applying a deck's own source changed it, on slide(s) {moved}")
print(f"   source {len(str(source_text or '')):,} chars · {len(before_deck.get('slides') or [])} slides unchanged")
call("DELETE", f"/presentations/{round_trip['id']}", expect=204)
for asset in (left, right):
    call("DELETE", f"/assets/{asset.get('id')}", expect=[204, 404])

# A deck whose every slide is skipped has no page to print. That is the person's
# to fix, so it is answered as such rather than as a server that broke.
skipped_deck = data_of(call("POST", "/presentations", {"title": f"전부 건너뜀 {RUN}", "prompt": "점검"}, expect=201))
call("PUT", f"/presentations/{skipped_deck['id']}/source",
     {"source": "# 부록\n@content\n!skip\n- 물어보면 보여 줄 표\n"}, expect=200)
status, refused = call("GET", f"/presentations/{skipped_deck['id']}/export?format=pdf", expect=409,
                       note="a deck with nothing to print must say so, not fail")
if (refused or {}).get("error", {}).get("code") != "presentation_has_no_printable_slides":
    failures.append(f"the refusal does not say what is wrong: {refused}")
call("GET", f"/presentations/{skipped_deck['id']}/export?format=pptx", raw=True, expect=200,
     note="the same deck still exports as pptx, where a hidden slide is still a slide")
call("DELETE", f"/presentations/{skipped_deck['id']}", expect=204)

# The handout: the same deck, with what the presenter meant to say under each
# slide. It is the notes that make it one, so an empty one is not a handout.
status, handout = call("GET", f"/presentations/{deck_id}/export?format=pdf&notes=true", raw=True, expect=200)
if not handout.startswith(b"%PDF-"):
    failures.append("the handout is not a PDF")
elif len(handout) <= len(pdf_bytes):
    failures.append(f"the handout ({len(handout):,}) carries no more than the deck ({len(pdf_bytes):,})")
else:
    print(f"   handout {len(handout):,} bytes")
# The notes are written the way every other line is, so the handout draws the
# words rather than the markup — and the link in them is one a reader follows.
notes_handout = data_of(call("POST", "/presentations", {"title": f"노트 마크업 {RUN}", "prompt": "실적"}, expect=201))
call("PUT", f"/presentations/{notes_handout['id']}/source", {"source":
     "# 근거\n@content\n- 매출이 늘었습니다\n"
     "!notes 근거는 [분기 보고서](https://reports.example.com/q3)에 있습니다\n!source 매출 | 사내 결산\n"}, expect=200)
status, printed = call("GET", f"/presentations/{notes_handout['id']}/export?format=pdf&notes=true", raw=True, expect=200)
if b"/URI (https://reports.example.com/q3)" not in printed:
    failures.append("the handout's notes carry no link a reader can follow")
call("DELETE", f"/presentations/{notes_handout['id']}", expect=204)

status, unsupported = call("GET", f"/presentations/{deck_id}/export?format=keynote", expect=422,
                           note="a format nothing can write must be refused with a reason")
if not (unsupported or {}).get("error", {}).get("code") == "unsupported_export_format":
    failures.append(f"the refusal does not say why: {unsupported}")

print("── canvas ──")
regions = data_of(call("GET", f"/presentations/{deck_id}/slides/2/regions", expect=200))
print("   regions on slide 2:", [r["slot"] for r in (regions or {}).get("regions", [])])
if not (regions or {}).get("regions"):
    failures.append("slide 2 reports no editable regions")
call("GET", f"/presentations/{deck_id}/slides/99/regions", expect=404)

print("── documents become decks ──")
sheet = "분기,매출\n1분기,1180\n2분기,1240\n3분기,1390\n"
imported = data_of(call("POST", "/presentations/import",
                        files={"file": (f"매출-{RUN}.csv", sheet, "text/csv")}, expect=201)) or {}
if (imported.get("slides") or 0) < 2:
    failures.append(f"a spreadsheet imported as {imported.get('slides')!r} slides")
document_id = ((imported.get("presentation") or {}).get("id")) or ""
# The source an author is shown is the deck. Applying it back must leave the
# deck as it was: a picture standing next to prose used to make the writer think
# the words were written, and one click in the editor deleted them.
if document_id:
    def visible(deck):
        seen = []
        for slide in (deck.get("slides") or []):
            content = slide.get("content")
            if isinstance(content, str):
                content = json.loads(content or "{}")
            seen.append(len(re.findall(r"[가-힣A-Za-z0-9]", json.dumps(content or {}, ensure_ascii=False))))
        return seen
    was = visible(data_of(call("GET", f"/presentations/{document_id}", expect=200)) or {})
    shown = (data_of(call("GET", f"/presentations/{document_id}/source", expect=200)) or {}).get("source", "")
    call("PUT", f"/presentations/{document_id}/source", {"source": shown}, expect=200)
    now = visible(data_of(call("GET", f"/presentations/{document_id}", expect=200)) or {})
    if len(was) != len(now):
        failures.append(f"applying a deck's own source changed it from {len(was)} slides to {len(now)}")
    else:
        for index, (before, after) in enumerate(zip(was, now), 1):
            if after < before * 0.6:
                failures.append(f"slide {index} lost most of its text to its own source: {before} → {after}")
# What the import had to say travels with the deck. It used to go into a toast
# that joined every sentence into one line and disappeared on the way to the
# editor — which has a panel for exactly this and was showing nothing.
if document_id and (imported.get("warnings") or []):
    carried = data_of(call("GET", f"/presentations/{document_id}", expect=200)) or {}
    kept = carried.get("generationNotes") or []
    if not kept:
        failures.append("an imported deck does not keep what the import said about the file")
    elif kept[0] != (imported.get("warnings") or [])[0]:
        failures.append(f"the deck kept {kept[0]!r} rather than what the import answered")
if document_id:
    written = (data_of(call("GET", f"/presentations/{document_id}/source", expect=200)) or {}).get("source", "")
    for wanted in ["::columns", "- 1분기 | 1180", f"!source 매출-{RUN}.csv"]:
        if wanted not in written:
            failures.append(f"the imported spreadsheet lost {wanted!r}")
    print(f"   spreadsheet -> {imported.get('slides')} slides, cited")
# A deck carried in from PowerPoint keeps what its runs said, not only what
# they spelled. An address is the one part of a link nobody can type again from
# looking at the slide, and it used to be dropped: the words came across and the
# link did not. The file is stripped of what Ptium leaves in its own exports, so
# the importer has to read the drawing the way it reads a stranger's deck.
# A link is the deck as it is shown. A slide the author marked skipped is not
# part of the show — PowerPoint hides it and the handout leaves it out — and it
# was being drawn to anyone holding the link.
print("── a link shows only what is in the show ──")
shared_deck = data_of(call("POST", "/presentations", {"title": f"공유 점검 {RUN}", "prompt": "점검",
                                                      "language": "ko"}, expect=201))
if shared_deck:
    call("PUT", f"/presentations/{shared_deck['id']}/source",
         {"source": "# 보이는 장\n- 공개해도 되는 요점\n\n# 숨긴 장\n- 아직 공유하지 않을 내용\n!skip\n"
                    "\n# 세 번째 장\n- 공개 요점\n"}, expect=200)
    made = data_of(call("POST", f"/presentations/{shared_deck['id']}/shares", {}, expect=201)) or {}
    token = (made.get("url") or "").rstrip("/").split("/")[-1]
    if token:
        seen = data_of(call("GET", f"/shared/{token}", expect=200)) or {}
        checks += 1
        if seen.get("slideCount") != 2:
            failures.append(f"a link with one skipped slide shows {seen.get('slideCount')} slides, want 2")
        checks += 1
        if any("숨긴" in (title or "") for title in (seen.get("titles") or [])):
            failures.append("a slide taken out of the show is named on the link")
        for position in (1, 2, 3):
            status, drawn = call("GET", f"/shared/{token}/preview.svg?slide={position}&width=400",
                                 raw=True, expect=200)
            checks += 1
            if "아직 공유하지 않을 내용" in drawn.decode(errors="replace"):
                failures.append(f"the link draws the skipped slide at position {position}")
        print(f"   slides on the link: {seen.get('slideCount')} of 3")
    call("DELETE", f"/presentations/{shared_deck['id']}", expect=204)

print("── a deck carried in keeps its links ──")
linked = data_of(call("POST", "/presentations", {"title": f"링크 왕복 {RUN}", "prompt": "점검",
                                                 "language": "ko"}, expect=201))
if linked:
    call("PUT", f"/presentations/{linked['id']}/source",
         {"source": "# 링크 슬라이드\n- 자료는 [계획서](https://example.com/plan)에, **굵게**도 있습니다\n"
                    "\n# 숨긴 장\n- 아직 공유하지 않을 내용\n!skip\n"
                    "\n# 근거\n- 요점\n!source 내부 자료 2026\n"
                    "!notes 자세한 것은 [문서](https://example.com/detail)에\n"}, expect=200)
    _, carried = call("GET", f"/presentations/{linked['id']}/export?format=pptx", raw=True, expect=200)
    stripped = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(carried)) as original:
        with zipfile.ZipFile(stripped, "w", zipfile.ZIP_DEFLATED) as out:
            for item in original.infolist():
                body = original.read(item.filename)
                if b"deckSource" in body:
                    if item.filename.endswith(".rels") or item.filename.endswith("[Content_Types].xml"):
                        body = re.sub(rb"<Relationship[^>]*deckSource[^>]*/>", b"", body)
                        body = re.sub(rb"<Override[^>]*ptiumSource[^>]*/>", b"", body)
                    else:
                        continue
                out.writestr(item, body)
    brought = data_of(call("POST", "/presentations/import",
                           files={"file": (f"남의덱-{RUN}.pptx", stripped.getvalue(),
                                           "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                           expect=201)) or {}
    brought_id = ((brought.get("presentation") or {}).get("id")) or ""
    if brought_id:
        read_back = (data_of(call("GET", f"/presentations/{brought_id}/source", expect=200)) or {}).get("source", "")
        for wanted in ["https://example.com/plan", "**굵게**", "!skip",
                       "!source 내부 자료 2026", "https://example.com/detail"]:
            checks += 1
            if wanted not in read_back:
                failures.append(f"a deck carried in from PowerPoint lost {wanted!r}")
        # A citation read back as a point would say "출처: …" in the body.
        checks += 1
        if "- 출처" in read_back:
            failures.append("a citation came back as a point rather than a citation")
        # A citation is written twice on the way out — drawn on the slide and
        # under the notes — and reading both back gave the deck the same source
        # as a citation and as a sentence in its notes, one more copy each trip.
        checks += 1
        note_lines = [line for line in read_back.splitlines() if line.startswith("!notes")]
        doubled = [line for line in note_lines if "내부 자료 2026" in line]
        if doubled:
            failures.append(f"a citation came back inside the notes as well: {doubled}")
        print(f"   link kept: {'https://example.com/plan' in read_back} · emphasis: {'**굵게**' in read_back}"
              f" · hidden: {'!skip' in read_back} · cited: {'!source' in read_back}")
        call("DELETE", f"/presentations/{brought_id}", expect=204)

    # Most real decks bold their title. A heading is a slot rather than prose —
    # the source has no way to bold part of one — so the marks come out as
    # characters: every real deck measured against this was called
    # "**목 차**" in the deck list and drew its own asterisks on the title slide.
    def rewritten(rule):
        made = io.BytesIO()
        with zipfile.ZipFile(io.BytesIO(stripped.getvalue())) as original:
            with zipfile.ZipFile(made, "w", zipfile.ZIP_DEFLATED) as out:
                for item in original.infolist():
                    body = original.read(item.filename)
                    if item.filename.startswith("ppt/slides/slide"):
                        body = rule(body)
                    out.writestr(item, body)
        return made.getvalue()

    emphatic = data_of(call("POST", "/presentations/import",
                            files={"file": (f"굵은제목-{RUN}.pptx",
                                            rewritten(lambda body: body.replace(b"<a:rPr ", b'<a:rPr b="1" ')),
                                            "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                            expect=201)) or {}
    emphatic_id = ((emphatic.get("presentation") or {}).get("id")) or ""
    if emphatic_id:
        checks += 1
        name = ((emphatic.get("presentation") or {}).get("title")) or ""
        if "*" in name:
            failures.append(f"a deck whose title is bold is called {name!r}")
        written = (data_of(call("GET", f"/presentations/{emphatic_id}/source", expect=200)) or {}).get("source", "")
        checks += 1
        if re.search(r"^# .*\*", written, re.M):
            failures.append(f"an imported heading carries its emphasis marks: "
                            f"{[line for line in written.splitlines() if line.startswith('#') and '*' in line][:2]}")
        print(f"   bold title read as: {name!r}")
        call("DELETE", f"/presentations/{emphatic_id}", expect=204)

    # A designer puts the section a slide belongs to in the corner of every
    # slide of that section, and the slide's own heading goes underneath. Read as
    # drawn, a twenty-two slide deck came in with fifteen slides all called
    # "3. 아이디어 구현" — an outline nobody can use. The slide keeps the name it
    # has when there is nothing else on it to be called.
    header_deck = data_of(call("POST", "/presentations", {"title": f"머리글 {RUN}", "prompt": "점검",
                                                          "language": "ko"}, expect=201))
    call("PUT", f"/presentations/{header_deck['id']}/source",
         {"source": "# 3. 아이디어 구현\n- 키워드 추출\n- 형태소 기준으로 뽑습니다\n"
                    "\n# 3. 아이디어 구현\n- 화자 분리\n- 두 사람을 나눕니다\n"
                    "\n# 3. 아이디어 구현\n- 감정 분석\n- 말투로 가늠합니다\n"
                    "\n# 4. 기대효과\n- 비용 절감\n"}, expect=200)
    _, headed = call("GET", f"/presentations/{header_deck['id']}/export?format=pptx", raw=True, expect=200)
    bare = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(headed)) as original:
        with zipfile.ZipFile(bare, "w", zipfile.ZIP_DEFLATED) as out:
            for item in original.infolist():
                body = original.read(item.filename)
                if b"deckSource" in body:
                    if item.filename.endswith(".rels") or item.filename.endswith("[Content_Types].xml"):
                        body = re.sub(rb"<Relationship[^>]*deckSource[^>]*/>", b"", body)
                        body = re.sub(rb"<Override[^>]*ptiumSource[^>]*/>", b"", body)
                    else:
                        continue
                out.writestr(item, body)
    marked = data_of(call("POST", "/presentations/import",
                          files={"file": (f"머리글-{RUN}.pptx", bare.getvalue(),
                                          "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                          expect=201)) or {}
    marked_id = ((marked.get("presentation") or {}).get("id")) or ""
    if marked_id:
        written = (data_of(call("GET", f"/presentations/{marked_id}/source", expect=200)) or {}).get("source", "")
        names = [line[2:].strip() for line in written.splitlines() if line.startswith("# ")]
        checks += 1
        if names.count("3. 아이디어 구현") > 1:
            failures.append(f"a section marker became the name of {names.count('3. 아이디어 구현')} slides: {names}")
        checks += 1
        if not {"키워드 추출", "화자 분리", "감정 분석"} <= set(names):
            failures.append(f"the slides' own headings are not their names: {names}")
        checks += 1
        if "머리글" not in " ".join(marked.get("warnings") or []):
            failures.append(f"a repeated header was dropped without saying so: {marked.get('warnings')}")
        print(f"   headers read as: {names}")
        call("DELETE", f"/presentations/{marked_id}", expect=204)
    call("DELETE", f"/presentations/{header_deck['id']}", expect=204)

    # A slide built in columns is read column by column. Every second real deck
    # has a roadmap drawn as stages side by side, and reading it down the page
    # interleaves the stages into one another — worse when the designer drew the
    # middle stage lower to make a zigzag, which came out as 1, 3, 2.
    INCH = 914400
    boxes = [(INCH // 2, INCH // 2, 11 * INCH, "향후 추진 계획"),
             (INCH // 2, 2 * INCH, 3 * INCH, "1단계 검증"),
             (INCH // 2, 3 * INCH, 3 * INCH, "1개월"),
             (9 * INCH // 2, 4 * INCH, 3 * INCH, "2단계 시범"),
             (9 * INCH // 2, 5 * INCH, 3 * INCH, "2개월"),
             (8 * INCH, 2 * INCH, 3 * INCH, "3단계 전환"),
             (8 * INCH, 3 * INCH, 3 * INCH, "3~6개월")]
    tree = ('<p:nvGrpSpPr><p:cNvPr id="1" name="Shape 1"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/>' +
            "".join(f'<p:sp><p:nvSpPr><p:cNvPr id="{index + 2}" name="Box {index + 2}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
                    f'<p:spPr><a:xfrm><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{INCH // 2}"/></a:xfrm></p:spPr>'
                    f'<p:txBody><a:bodyPr/><a:p><a:r><a:rPr lang="ko-KR"/><a:t>{text}</a:t></a:r></a:p></p:txBody></p:sp>'
                    for index, (left, top, width, text) in enumerate(boxes)))
    columns = data_of(call("POST", "/presentations/import",
                           files={"file": (f"단쌓기-{RUN}.pptx",
                                           rewritten(lambda body: re.sub(rb"<p:spTree>.*?</p:spTree>",
                                                                         ("<p:spTree>" + tree + "</p:spTree>").encode(),
                                                                         body, count=1, flags=re.S)),
                                           "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                           expect=201)) or {}
    columns_id = ((columns.get("presentation") or {}).get("id")) or ""
    if columns_id:
        written = (data_of(call("GET", f"/presentations/{columns_id}/source", expect=200)) or {}).get("source", "")
        read_lines = [line.lstrip("# -").strip() for line in written.splitlines()
                      if line.startswith("# ") or line.startswith("- ")]
        checks += 1
        wanted = ["향후 추진 계획", "1단계 검증", "1개월", "2단계 시범", "2개월", "3단계 전환", "3~6개월"]
        if read_lines[:len(wanted)] != wanted:
            failures.append(f"a slide built in columns was read as {read_lines[:len(wanted)]}")
        print(f"   columns read as: {' · '.join(read_lines[:4])}")
        call("DELETE", f"/presentations/{columns_id}", expect=204)

    # A deck exported as one picture per slide — what several deck generators
    # produce — imports as slides called "3번 슬라이드" with nothing on them. That
    # is what the file holds; saying nothing about it reads as a failed import.
    wordless = data_of(call("POST", "/presentations/import",
                            files={"file": (f"그림만-{RUN}.pptx",
                                            rewritten(lambda body: re.sub(rb"<a:t>.*?</a:t>", b"<a:t></a:t>", body, flags=re.S)),
                                            "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                            expect=201)) or {}
    wordless_id = ((wordless.get("presentation") or {}).get("id")) or ""
    if wordless_id:
        checks += 1
        said = " | ".join(wordless.get("warnings") or [])
        if "글자가 없" not in said:
            failures.append(f"a deck of slides with nothing to read was imported saying {said!r}")
        print(f"   wordless deck told: {said[:60]!r}")
        call("DELETE", f"/presentations/{wordless_id}", expect=204)
    call("DELETE", f"/presentations/{linked['id']}", expect=204)

call("POST", "/presentations/import", files={"file": ("보고서.pdf", b"%PDF-1.7", "application/pdf")}, expect=422,
     note="a PDF this cannot open must say so rather than import an empty deck")

# A .pptx that will not open is a presentation with something wrong with it —
# usually document security. Sent to the document reader it was answered
# "Ptium reads .pptx presentations and …", about a .pptx.
locked = b"SCDS" + bytes(512)
status, refused = call("POST", "/presentations/import",
                       files={"file": (f"잠긴 파일-{RUN}.pptx", locked,
                                       "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                       expect=422, note="a locked presentation must be named for what it is")
said = ((refused or {}).get("error") or {}).get("message", "") if isinstance(refused, dict) else str(refused)
checks += 1
if "DRM" not in said:
    failures.append(f"a document-security wrapper was refused as {said[:120]!r}")
checks += 1
if "reads" in json.dumps(refused, ensure_ascii=False):
    failures.append(f"a .pptx was told the product does not read .pptx: {said[:120]!r}")
print("   locked .pptx told:", said[:56])

print("── a PDF becomes slides ──")
# A PDF is what a company sends things in, and it is also the format that keeps
# the least: glyphs at coordinates, no headings and no points. What can be
# recovered is each page's own lines; what cannot be recovered is said out loud.
report = pdf_of([
    ["국가법령정보센터", "- 1 -", f"2026 상반기 실적 {RUN}", "매출은 1,240억입니다.", "영업이익은 210억입니다."],
    ["국가법령정보센터", "- 2 -", "하반기 계획", "신규 채널 두 곳을 엽니다."],
    # More lines than a slide holds: the rest belongs in the notes, with the
    # document's own punctuation intact — a bar in a report is punctuation, not
    # a separator.
    ["국가법령정보센터", "- 3 -", "위험 요인"] + [f"단가 | 수수료 | 거래세 {n}" for n in range(1, 9)],
])
brought_pdf = data_of(call("POST", "/presentations/import",
                           files={"file": (f"실적 보고-{RUN}.pdf", report, "application/pdf")},
                           expect=201)) or {}
pdf_id = ((brought_pdf.get("presentation") or {}).get("id")) or ""
if pdf_id:
    source = (data_of(call("GET", f"/presentations/{pdf_id}/source", expect=200)) or {}).get("source", "")
    checks += 1
    if "# 하반기 계획" not in source or "# 위험 요인" not in source:
        failures.append(f"a PDF's pages did not become slides:\n{source[:400]}")
    checks += 1
    if "!source 실적 보고-" not in source or "쪽" not in source:
        failures.append(f"a slide from a PDF does not cite the page it came from:\n{source[:400]}")
    checks += 1
    if "국가법령정보센터" in source:
        failures.append(f"the header repeated on every page became slide text:\n{source[:400]}")
    checks += 1
    if "- - 2 -" in source or "# - 2 -" in source:
        failures.append(f"a page number became slide text:\n{source[:400]}")
    checks += 1
    if "!notes 단가 | 수수료 | 거래세 7" not in source:
        failures.append(f"what did not fit on the slide is in no note, or its punctuation was rewritten:\n{source[-600:]}")
    print("   PDF read as:", " · ".join(line[2:] for line in source.splitlines() if line.startswith("# "))[:70])
    call("DELETE", f"/presentations/{pdf_id}", expect=204)

# A document is usually named after what it is about, so the cover's title and
# the line under it — the file it came from — say the same thing. The product's
# own measurement calls that "the same point twice" on a deck nobody has edited.
named = f"협업 도구 소개 {RUN}"
same = data_of(call("POST", "/presentations/import",
                    files={"file": (f"{named}.pdf",
                                    pdf_of([[named, "여는 문장입니다."], ["도입 효과", "둘째 쪽입니다."]]),
                                    "application/pdf")},
                    expect=201)) or {}
same_id = ((same.get("presentation") or {}).get("id")) or ""
if same_id:
    measured = data_of(call("GET", f"/presentations/{same_id}/inspect", expect=200)) or {}
    repeated = [f for f in (measured.get("findings") or [])
                if f.get("kind") == "repeat" and f.get("slide") == 1]
    checks += 1
    if repeated:
        failures.append(f"a freshly imported cover repeats itself: {repeated}")
    source = (data_of(call("GET", f"/presentations/{same_id}/source", expect=200)) or {}).get("source", "")
    checks += 1
    if f"# {named}" not in source:
        failures.append(f"the cover lost its title:\n{source[:200]}")
    print(f"   cover of a file named after itself: {len(repeated)} repeat finding(s)")
    call("DELETE", f"/presentations/{same_id}", expect=204)

# A scan has no text in it. An import that answers with an empty deck teaches
# people the product does not work; one that says what is wrong tells them what
# to upload instead.
status, refused = call("POST", "/presentations/import",
                       files={"file": (f"스캔본-{RUN}.pdf", pdf_of([["x"]], image_only=True), "application/pdf")},
                       expect=422, note="a scanned PDF must say why it cannot be read")
# The reason has to be in the message, not merely somewhere in the reply: the
# name of the file the caller uploaded is echoed back in it, and a check that
# searched the whole body would pass on the word in the filename.
said = ((refused or {}).get("error") or {}).get("message", "") if isinstance(refused, dict) else str(refused)
checks += 1
if "스캔" not in said or "PDF" not in said:
    failures.append(f"a scanned PDF was refused without saying why: {said[:200]!r}")
print("   scan refused with:", said[:70])

# What a file is called decides nothing about what it is. A browser sends
# "image/png" for whatever the person picked, and a picture that is not a
# picture is found out by PowerPoint, which calls the exported deck damaged.
for what, bytes_of in {"a line of text": b"nope not an image at all",
                       "a web page": b"<!doctype html><html><body>hi</body></html>",
                       "a PDF": b"%PDF-1.7\n1 0 obj\n<</Type /Catalog>>"}.items():
    status, refused = call("POST", "/assets",
                           files={"file": (f"{RUN}-not-an-image.png", bytes_of, "image/png")},
                           expect=422, note=f"{what} named .png must be refused")
    if status in (200, 201):
        stored = data_of((status, refused)) or {}
        if stored.get("id"):
            call("DELETE", f"/assets/{stored['id']}", expect=[204, 404])
checks += 1
real = data_of(call("POST", "/assets", files={"file": (f"{RUN}-real.png", png(), "image/png")}, expect=201)) or {}
if not real.get("id"):
    failures.append("a real PNG was refused along with the impostors")
else:
    call("DELETE", f"/assets/{real['id']}", expect=[204, 404])
print("   the picture library takes pictures and nothing else")

# The wildcards in a search belong to the query, not to the reader. "%"+text+"%"
# handed what somebody typed straight to LIKE, so searching for "50%" matched
# "500만원" and a name with an underscore could not be searched for exactly.
wild = {}
for name in [f"달성률 50% 보고 {RUN}", f"달성률 500만원 보고 {RUN}",
             f"metadata_demo 정리 {RUN}", f"metadataXdemo 정리 {RUN}"]:
    made_deck = data_of(call("POST", "/presentations", {"title": name, "prompt": "검색 점검",
                                                        "language": "ko"}, expect=201)) or {}
    wild[name] = made_deck.get("id")
for typed, wanted in ((f"50% 보고 {RUN}", f"달성률 50% 보고 {RUN}"),
                      (f"metadata_demo 정리 {RUN}", f"metadata_demo 정리 {RUN}")):
    found = data_of(call("GET", f"/presentations?q={urllib.parse.quote(typed)}&limit=10", expect=200)) or []
    titles = [row.get("title") for row in found]
    checks += 1
    if titles != [wanted]:
        failures.append(f"searching {typed!r} found {titles}")
print("   searching for a % or a _ finds the name that has one")
# Half of this API called the filter "q" and half "search", and the name that
# was not read was ignored in silence — a client searching templates with "q"
# was handed every template.
for where, term in (("/presentations", f"50% 보고 {RUN}"), ("/templates", "Rail"),
                    ("/assets", "e2e"), ("/snippets", "e2e")):
    totals = []
    for name in ("q", "search"):
        page = call("GET", f"{where}?{name}={urllib.parse.quote(term)}&limit=5", expect=200)
        totals.append(((page[1] or {}).get("meta") or {}).get("total"))
    checks += 1
    if totals[0] != totals[1]:
        failures.append(f"{where} answers q={totals[0]} and search={totals[1]} for the same question")
print("   q and search ask the same question at every door")
for one in wild.values():
    if one:
        call("DELETE", f"/presentations/{one}?permanent=true", expect=[204, 404])

# The recycle bin could only be cleared one deck at a time, which is no way to
# clear a thousand. What it takes is what is already in it — and nothing else.
kept = data_of(call("POST", "/presentations", {"title": f"살아 있어야 하는 덱 {RUN}", "prompt": "점검",
                                               "language": "ko"}, expect=201)) or {}
doomed = data_of(call("POST", "/presentations", {"title": f"지워질 덱 {RUN}", "prompt": "점검",
                                                 "language": "ko"}, expect=201)) or {}
call("DELETE", f"/presentations/{doomed['id']}", expect=204)
emptied = data_of(call("DELETE", "/presentations/trash", expect=200)) or {}
checks += 1
if not isinstance(emptied.get("deleted"), int) or emptied.get("deleted") < 1:
    failures.append(f"emptying the trash reported {emptied!r}")
checks += 1
if call("GET", f"/presentations/{kept['id']}")[0] != 200:
    failures.append("emptying the trash took a deck that was not in it")
checks += 1
if call("GET", f"/presentations/{doomed['id']}")[0] != 404:
    failures.append("a deck in the trash survived emptying it")
checks += 1
still = data_of(call("GET", "/presentations?deleted=true&limit=5", expect=200))
if still:
    failures.append(f"the trash still holds {len(still)} deck(s) after being emptied")
print("   trash emptied:", emptied.get("deleted"), "· the live deck is untouched")
call("DELETE", f"/presentations/{kept['id']}?permanent=true", expect=[204, 404])

# Somebody who wrote 62 slides has one of the two numbers and needs the other:
# "more than this deployment allows" leaves them counting their own slides.
too_many = data_of(call("POST", "/presentations", {"title": f"경계 {RUN}", "prompt": "점검",
                                                   "language": "ko"}, expect=201)) or {}
if too_many.get("id"):
    long_source = "".join(f"# {n}번째 장\n@content\n- 요점 {n}\n\n" for n in range(1, 62))
    status, refused = call("PUT", f"/presentations/{too_many['id']}/source", {"source": long_source},
                           expect=422, note="a deck past the limit must say what the limit is")
    said = ((refused or {}).get("error") or {}).get("message", "") if isinstance(refused, dict) else ""
    checks += 1
    if "61" not in said or not any(str(limit) in said for limit in (50, 60, 100, 200)):
        failures.append(f"a deck past the limit was refused with {said!r}")
    print("   past the limit:", said[:70])
    call("DELETE", f"/presentations/{too_many['id']}?permanent=true", expect=[204, 404])

print("── a link someone without an account can open ──")
shared_deck = data_of(call("POST", "/presentations", {"title": f"공유 점검 {RUN}", "prompt": "공유",
                                                      "requestedSlideCount": 3, "language": "ko"}, expect=201)) or {}
call("PUT", f"/presentations/{shared_deck.get('id')}/source",
     {"source": f"# 공유 점검 {RUN}\n@cover\n> 링크로 여는 덱\n\n# 현황\n@content\n- 첫 줄\n"
                "- 근거는 [자료](https://example.com/근거)에 있습니다\n"})
made = data_of(call("POST", f"/presentations/{shared_deck.get('id')}/shares", {"label": "임원 검토", "days": 7}, expect=201)) or {}
link = made.get("url", "")
token = link.rsplit("/", 1)[-1] if link else ""
if not token or "/view/" not in link:
    failures.append(f"a share was made without a link to hand anyone: {made!r}")
# Opened by a stranger: no dev secret, no session, no key.
opened = urllib.request.Request(BASE + f"/shared/{token}")
try:
    with urllib.request.urlopen(opened) as response:
        seen = json.loads(response.read())["data"]
except urllib.error.HTTPError as error:
    seen = {}
    failures.append(f"a shared link did not open for someone without an account: {error.code}")
if seen.get("slideCount") != 2 or f"공유 점검 {RUN}" not in (seen.get("title") or ""):
    failures.append(f"the link showed something other than the deck: {seen!r}")
try:
    with urllib.request.urlopen(urllib.request.Request(BASE + f"/shared/{token}/preview.svg?slide=2&width=600")) as response:
        drawn = response.read().decode()
    if "<svg" not in drawn or "자료" not in drawn:
        failures.append("a shared link drew something other than the slide")
    # A deck sent out for review is where links are read: the page draws the
    # slide itself, so the drawing has to carry them.
    if '<a href="https://example.com/' not in drawn:
        failures.append("a shared slide draws its links as words nobody can follow")
except urllib.error.HTTPError as error:
    failures.append(f"a shared slide did not draw: {error.code}")
# The link carries nothing else: the deck's source is not readable through it.
for path in [f"/presentations/{shared_deck.get('id')}", f"/presentations/{shared_deck.get('id')}/source"]:
    try:
        urllib.request.urlopen(urllib.request.Request(BASE + path))
        failures.append(f"{path} answered someone with no credentials")
    except urllib.error.HTTPError as error:
        if error.code not in (401, 403):
            failures.append(f"{path} answered {error.code} to someone with no credentials")
call("DELETE", f"/presentations/{shared_deck.get('id')}/shares/{made.get('id')}", expect=204)
try:
    urllib.request.urlopen(urllib.request.Request(BASE + f"/shared/{token}"))
    failures.append("a revoked link still opens the deck")
except urllib.error.HTTPError as error:
    if error.code != 410:
        failures.append(f"a revoked link answered {error.code}, expected 410")
# Looking is half of a review. The other half is saying what is wrong with
# slide 4, which is what the link is for.
opened_again = data_of(call("POST", f"/presentations/{shared_deck.get('id')}/shares", {"label": "의견", "days": 0}, expect=201)) or {}
review_token = (opened_again.get("url") or "").rsplit("/", 1)[-1]
review_slides = []
try:
    with urllib.request.urlopen(urllib.request.Request(BASE + f"/shared/{review_token}")) as response:
        review_slides = (json.loads(response.read())["data"] or {}).get("slides") or []
except urllib.error.HTTPError as error:
    failures.append(f"the review link did not open: {error.code}")
if review_slides:
    said = json.dumps({"slideId": review_slides[1]["id"], "author": "박검토",
                       "body": "이 숫자는 지난 분기 기준입니다."}).encode()
    request = urllib.request.Request(BASE + f"/shared/{review_token}/comments", data=said,
                                     headers={"Content-Type": "application/json"}, method="POST")
    try:
        with urllib.request.urlopen(request) as response:
            left = json.loads(response.read())["data"]
        if left.get("slideId") != review_slides[1]["id"]:
            failures.append(f"the remark lost the slide it was about: {left!r}")
    except urllib.error.HTTPError as error:
        failures.append(f"a stranger could not leave a remark: {error.code}")
    mine = data_of(call("GET", f"/presentations/{shared_deck.get('id')}/comments")) or []
    if not any(row.get("author") == "박검토" for row in mine):
        failures.append("the owner cannot see what the reviewer said")
    else:
        remark = [row for row in mine if row.get("author") == "박검토"][0]
        call("POST", f"/presentations/{shared_deck.get('id')}/comments/{remark['id']}/resolve", {"resolved": True}, expect=204)
        after = data_of(call("GET", f"/presentations/{shared_deck.get('id')}/comments")) or []
        if not any(row["id"] == remark["id"] and row.get("resolvedAt") for row in after):
            failures.append("a remark marked dealt with did not stay dealt with")
    # A remark about someone else's slide is not a way into another deck.
    stray = json.dumps({"slideId": "11111111-1111-1111-1111-111111111111", "body": "x"}).encode()
    request = urllib.request.Request(BASE + f"/shared/{review_token}/comments", data=stray,
                                     headers={"Content-Type": "application/json"}, method="POST")
    try:
        urllib.request.urlopen(request)
        failures.append("a remark was accepted for a slide from another deck")
    except urllib.error.HTTPError as error:
        if error.code != 422:
            failures.append(f"a stray slide id answered {error.code}, expected 422")
print("   a reviewer left a remark on slide 2 and the owner marked it dealt with")

print(f"   opened by a stranger, {seen.get('slideCount')} slides, closed on revoke")

print("── the library comes first ──")
registered = data_of(call("POST", "/snippets", {
    "name": f"회사 소개 {RUN}",
    "source": f"# 회사 소개 {RUN}\n@content\n> 2003년 설립\n- 임직원 1,240명\n- 매출 8,200억\n",
    "tags": ["표준"],
}, expect=201)) or {}
queued = data_of(call("POST", "/presentations/generate", {
    "prompt": f"회사 소개 {RUN}와 2026년 사업 계획을 임원에게 보고", "requestedSlideCount": 6, "language": "ko",
}, expect=202)) or {}
library_deck = queued.get("id", "")
for _ in range(40):
    state = data_of(call("GET", f"/presentations/{library_deck}")) or {}
    if state.get("status") not in ("queued", "generating"):
        break
    time.sleep(1)
written = (data_of(call("GET", f"/presentations/{library_deck}/source", expect=200)) or {}).get("source", "")
if "임직원 1,240명" not in written:
    failures.append("generation wrote its own version of a slide the library already had")
# Asked for by id: a workspace with more than a page of favourites pushes a new
# slide off the front of the list, and reading only the front page reported the
# product broken when it was the listing that had moved on.
used = data_of(call("GET", f"/snippets/{registered.get('id')}", expect=200)) or {}
if used.get("useCount", 0) < 1:
    failures.append(f"the registered slide was not counted as used: {used!r}")
else:
    print(f"   the deck took the registered slide (used {used.get('useCount')}x)")

print("── commanding the deck ──")
plan = data_of(call("POST", f"/presentations/{deck_id}/command", {"text": "2번과 3번 합쳐줘", "dryRun": True}, expect=200)) or {}
if not (plan.get("plan") or []):
    failures.append(f"a command was understood as nothing: {plan!r}")
elif plan.get("applied"):
    failures.append("a dry run changed the deck")
else:
    print(f"   plan: {plan['plan'][0].get('reason')} ({plan.get('slides')} -> {plan.get('slidesAfter')})")
call("POST", f"/presentations/{deck_id}/command", {"text": "이 덱 어때?"}, expect=422,
     note="a sentence that names no edit must say so rather than guessing")

print("── quality score ──")
measured = data_of(call("GET", f"/presentations/{deck_id}/inspect", expect=200)) or {}
score = measured.get("score") or {}
if not (0 <= int(score.get("total", -1)) <= 100):
    failures.append(f"the deck's quality score is {score.get('total')!r}")
if len(score.get("slides") or []) != measured.get("slides"):
    failures.append(f"the score covers {len(score.get('slides') or [])} of {measured.get('slides')} slides")
if {entry.get("key") for entry in (score.get("dimensions") or [])} != {
        "readability", "structure", "visual", "accessibility", "evidence"}:
    failures.append(f"the score's dimensions are {score.get('dimensions')!r}")
print(f"   score {score.get('total')} weakest slide {score.get('weakest')}")

# Three answers, and they used to be one. A sentence naming a slide the deck
# does not have was understood; so was one asking for what the deck already is.
# Telling their authors that nothing they said made sense sends them looking for
# better words rather than for the right number.
already = data_of(call("GET", f"/presentations/{deck_id}", expect=200)) or {}
count = len(already.get("slides") or [])
if count:
    status, answered = call("POST", f"/presentations/{deck_id}/command",
                            {"text": f"{count}장으로 줄여줘", "dryRun": True}, expect=200)
    said = data_of((status, answered)) or {}
    checks += 1
    if said.get("plan"):
        failures.append(f"a deck already {count} slides long was told to trim: {said.get('plan')}")
    checks += 1
    if not any("이미" in note for note in (said.get("notes") or [])):
        failures.append(f"asking for what the deck already is answered {said.get('notes')!r}")
    status, refused = call("POST", f"/presentations/{deck_id}/command",
                           {"text": f"{count + 40}번 슬라이드 삭제", "dryRun": True}, expect=422)
    checks += 1
    if "없습니다" not in json.dumps(refused, ensure_ascii=False):
        failures.append(f"naming a slide the deck has not answered {refused!r}")
    checks += 1
    if "command_out_of_range" not in json.dumps(refused, ensure_ascii=False):
        failures.append(f"a slide out of range was reported as not understood: {refused!r}")
    print(f"   already {count} slides: told so · slide {count + 40}: named")

print("── images ──")
asset = data_of(call("POST", "/assets", files={"file": (f"logo-{RUN}.png", png(rgb=(int(RUN[-3:]) % 255, 120, 60)), "image/png")}, expect=201))
asset_id = asset["id"]
same = call("POST", "/assets", files={"file": (f"logo-copy-{RUN}.png", png(rgb=(int(RUN[-3:]) % 255, 120, 60)), "image/png")}, expect=200)
if not (data_of(same) or {}).get("reused"):
    failures.append("uploading identical bytes did not report a reuse")
call("POST", "/assets", files={"file": ("note.txt", b"not an image", "text/plain")}, expect=422)
call("PATCH", f"/assets/{asset_id}", {"tags": ["로고", "E2E"]}, expect=200)
call("PUT", f"/assets/{asset_id}/favorite", {"favorite": True}, expect=200)
listed = data_of(call("GET", "/assets?favorite=true&sort=used", expect=200)) or []
if not any(a["id"] == asset_id for a in listed):
    failures.append("a starred image is missing from the starred list")
call("GET", "/assets/tags", expect=200)
call("GET", f"/assets/{asset_id}", raw=True, expect=200)
call("PATCH", f"/assets/{asset_id}", {"name": ""}, expect=422)

print("── images inside a deck ──")
with_image = source + f"\n# 그림\n@picture\n::image logo-{RUN}.png\n"
call("PUT", f"/presentations/{deck_id}/source", {"source": with_image}, expect=200)
used = [a for a in (data_of(call("GET", "/assets", expect=200)) or []) if a["id"] == asset_id]
if not used or used[0]["deckCount"] != 1:
    failures.append(f"the image should be counted in one deck: {used}")
status, pptx = call("GET", f"/presentations/{deck_id}/export?format=pptx", raw=True, expect=200)
media = [n for n in zipfile.ZipFile(io.BytesIO(pptx)).namelist() if n.startswith("ppt/media/")]
if not media:
    failures.append("the exported deck carries no image although one is placed")

print("── saved slides ──")
snippet = data_of(call("POST", "/snippets", {"presentationId": deck_id, "slide": 2, "tags": ["표준"]}, expect=201))
snippet_id = snippet["id"]
print("   saved:", snippet["name"])
call("GET", "/snippets?sort=lastUsed", expect=200)
call("GET", "/snippets/tags", expect=200)
rendered = data_of(call("POST", f"/snippets/{snippet_id}/render", {"presentationId": deck_id}, expect=200))
if not (rendered or {}).get("slide"):
    failures.append("rendering a saved slide returned nothing")
call("GET", f"/snippets/{snippet_id}/preview.svg?presentationId={deck_id}&width=320", raw=True, expect=200)
call("PUT", f"/snippets/{snippet_id}/favorite", {"favorite": True}, expect=200)
call("PATCH", f"/snippets/{snippet_id}", {"name": f"회사 표준 지표 {RUN}"}, expect=200)
second = data_of(call("POST", "/snippets", {"source": "# 다른 조각\n- 한 줄\n", "name": f"다른 조각 {RUN}"}, expect=201))
call("PATCH", f"/snippets/{second['id']}", {"name": f"회사 표준 지표 {RUN}"}, expect=409, note="a taken name must conflict")
call("DELETE", f"/snippets/{second['id']}", expect=204)
call("GET", f"/snippets/{second['id']}", expect=404)

print("── grids ──")
call("GET", "/grids", expect=200)
call("PUT", f"/grids/e2e-raci-{RUN}", {"name": f"e2e-raci-{RUN}", "title": "E2E RACI",
     "columns": [{"label": "역할"}, {"label": "책임"}],
     "values": {"A": {"role": "accent1", "label": "승인"}, "R": {"role": "positive", "label": "실행"}},
     "order": ["R", "A"], "legend": True}, expect=[200, 201])
call("DELETE", f"/grids/e2e-raci-{RUN}", expect=[204, 200])

# An OIDC client secret can be registered in the workspace rather than in a
# deployment manifest. It is stored encrypted, never read back, and a secret
# with no client to belong to is refused before it can stop the next rollout.
print("── the OIDC client secret ──")
def oidc_secret_row():
    rows = data_of(call("GET", "/admin/settings", expect=200)) or []
    rows = rows.get("settings") if isinstance(rows, dict) else rows
    return next((r for r in rows if r.get("key") == "auth.oidc.client_secret"), {})

call("PUT", "/admin/settings", {"settings": [{"key": "auth.oidc.client_id", "value": ""},
                                             {"key": "auth.oidc.client_secret", "value": "orphan-secret"}]},
     expect=422, note="a client secret with no client ID must be refused")
call("PUT", "/admin/settings", {"settings": [{"key": "auth.oidc.client_id", "value": f"ptium-web-{RUN}"}]}, expect=200)
call("PUT", "/admin/settings", {"settings": [{"key": "auth.oidc.client_secret", "value": f"secret-{RUN}"}]}, expect=200)
if not oidc_secret_row().get("configured"):
    failures.append("a saved OIDC client secret is not reported as configured")
status, body = call("GET", "/admin/settings", raw=True, expect=200)
if f"secret-{RUN}".encode() in body:
    failures.append("the OIDC client secret is echoed back in the settings response")
call("PUT", "/admin/settings", {"settings": [{"key": "auth.oidc.client_secret", "value": " "}]}, expect=200)
if oidc_secret_row().get("configured"):
    failures.append("clearing the OIDC client secret left it configured")
call("PUT", "/admin/settings", {"settings": [{"key": "auth.oidc.client_id", "value": ""}]}, expect=200)

print("── revisions, duplicate, trash ──")
revisions = data_of(call("GET", f"/presentations/{deck_id}/revisions", expect=200)) or []
print("   revisions:", len(revisions))
if revisions:
    # What changed since a version — the question version history raises.
    changes = (data_of(call("GET", f"/presentations/{deck_id}/revisions/{revisions[0]['id']}/changes",
                            expect=200)) or {}).get("changes")
    if changes is None:
        failures.append("a revision cannot say what changed since it")
    call("POST", f"/presentations/{deck_id}/revisions/{revisions[-1]['id']}/restore", {}, expect=200)

# Restoring a checkpoint has to give the deck back. The call was checked for a
# 200 and nothing else, so nothing here would have noticed a restore that
# returned a different deck.
call("PUT", f"/presentations/{deck_id}/source",
     {"source": f"# 복원 점검 {RUN}\n@cover\n> 되돌리기 전\n\n# 요점\n@content\n- 첫 줄\n- 둘째 줄\n"}, expect=200)
before = ((data_of(call("GET", f"/presentations/{deck_id}/source", expect=200)) or {}).get("source") or "").strip()
call("PUT", f"/presentations/{deck_id}/source",
     {"source": f"# 덮어쓴 판 {RUN}\n@cover\n> 이 판은 되돌려집니다\n"}, expect=200)
# A checkpoint records the deck as it was before the change, so the newest one
# holds what the overwrite replaced.
newest = (data_of(call("GET", f"/presentations/{deck_id}/revisions", expect=200)) or [])
if not newest:
    failures.append("overwriting a deck kept no checkpoint of what it replaced")
else:
    # Which checkpoint this is matters as much as what it gives back: the
    # overwrite is what made it, so anything else at the front of the list means
    # the checkpoint the restore needs was never written.
    marker = newest[0]
    call("POST", f"/presentations/{deck_id}/revisions/{marker['id']}/restore", {}, expect=200)
    back = ((data_of(call("GET", f"/presentations/{deck_id}/source", expect=200)) or {}).get("source") or "").strip()
    if back != before:
        failures.append(
            "a restored checkpoint did not give the deck back — newest checkpoint was "
            f"{marker.get('reason')} at version {marker.get('version')} ({marker.get('createdAt')}); "
            f"the deck now starts {back.splitlines()[:1]} and should start {before.splitlines()[:1]}; "
            f"checkpoints: {[(r.get('reason'), r.get('version')) for r in newest[:5]]}")
        print("   lost:", [l for l in before.split("\n") if l not in back.split("\n")][:4])
        print("   gained:", [l for l in back.split("\n") if l not in before.split("\n")][:4])

copy = data_of(call("POST", f"/presentations/{deck_id}/duplicate", {}, expect=201))
call("DELETE", f"/presentations/{copy['id']}", expect=204)
trashed = data_of(call("GET", "/presentations?deleted=true", expect=200)) or []
if not any(p["id"] == copy["id"] for p in trashed):
    failures.append("a deleted deck is not in the recycle bin")
call("POST", f"/presentations/{copy['id']}/restore", {}, expect=200)
call("DELETE", f"/presentations/{copy['id']}", expect=204)
call("DELETE", f"/presentations/{copy['id']}?permanent=true", expect=[204, 200, 404])

# A brief that is mostly numbers earns a slide that draws them. Without a model
# this deck used to come back as bullet lists: of the decks this workspace has
# generated, 461 with a chart were written by hand and three by generation.
# The product's central promise: the template is the design, and a deck that was
# written into one can be written into another. A deck rebinds each slide to the
# layout of the same role, and what it draws changes with it.
# A box off the network has one disk, and when it fills the failures arrive as
# whatever the layer underneath happens to raise. The administrator's screens
# said nothing about it.
# Setting a provider and learning whether it answers were two different days:
# the only way to find out was to generate a deck and watch it fail. The
# administrator's own screen said "연결됨" whichever was true.
print("── whether the model host answers ──")
def setting_rows():
    """The settings, whichever shape this deployment answers with."""
    answered = data_of(call("GET", "/admin/settings", expect=200))
    if isinstance(answered, dict):
        return answered.get("settings") or answered.get("items") or []
    return answered or []


PROVIDER_BASE_URL = next((row.get("value") for row in setting_rows()
                          if row.get("key") == "ai.base_url"), "https://api.openai.com/v1")
was_provider = next((row.get("value") for row in setting_rows() if row.get("key") == "ai.provider"), None)
checked = data_of(call("POST", "/admin/provider-check", {}, expect=200)) or {}
checks += 1
if checked.get("provider") is None:
    failures.append("the provider check does not say which provider it asked")
checks += 1
if checked.get("provider") == "fallback" and checked.get("reachable"):
    failures.append("a deployment with no provider reports one answering")
# Pointed at a port nothing listens on, it says so rather than claiming health.
call("PUT", "/admin/settings", {"values": {"ai.provider": "openai-compatible",
                                           "ai.base_url": "http://127.0.0.1:9/v1"}}, expect=200)
missing = data_of(call("POST", "/admin/provider-check", {}, expect=200)) or {}
checks += 1
if missing.get("reachable"):
    failures.append("a provider pointed at nothing is reported as reachable")
checks += 1
if not str(missing.get("detail") or "").strip():
    failures.append("a provider that did not answer gives no reason why")
# A reading of a thing that has changed is not a reading of it: repointing the
# provider must not leave a screen reporting the host it used to be. The screen
# is asked straight after the change with nothing in between — an explicit check
# refreshes the reading, and would hide this.
checks += 1
if missing.get("baseUrl") != "http://127.0.0.1:9/v1":
    failures.append(f"after repointing the provider, the check still reads {missing.get('baseUrl')!r}")
call("PUT", "/admin/settings", {"values": {"ai.base_url": PROVIDER_BASE_URL}}, expect=200)
seen_by_a_screen = data_of(call("GET", "/admin/provider-check", expect=200)) or {}
checks += 1
if seen_by_a_screen.get("baseUrl") != PROVIDER_BASE_URL:
    failures.append("a screen reads the provider as it was before the change: "
                    f"{seen_by_a_screen.get('baseUrl')!r}, want {PROVIDER_BASE_URL!r}")
call("PUT", "/admin/settings", {"values": {"ai.provider": was_provider or "fallback",
                                           "ai.base_url": PROVIDER_BASE_URL}}, expect=200)
call("POST", "/admin/provider-check", {}, expect=403,
     headers={"X-Ptium-Dev-Secret": SECRET, "Authorization": f"Bearer dev:plain-{RUN}@ptium.local:user"},
     note="asking somebody's model host is the administrator's")
# A screen showing the state is not a person asking for it. A dashboard that
# writes an audit entry every time it opens fills the trail with itself.
def provider_check_entries():
    answered = call("GET", "/admin/audit?action=settings.provider_check&days=1&limit=1", expect=200)[1]
    return int(((answered or {}).get("meta") or {}).get("total") or 0)


trail_before = provider_check_entries()
for _ in range(3):
    call("GET", "/admin/provider-check", expect=200)
checks += 1
if provider_check_entries() != trail_before:
    failures.append("a screen reading the provider's state writes itself into the audit trail")
call("POST", "/admin/provider-check", {}, expect=200)
checks += 1
if provider_check_entries() != trail_before + 1:
    failures.append("a person asking the provider is not written down")
print(f"   asked: provider={checked.get('provider')!r} reachable={checked.get('reachable')} "
      f"{checked.get('milliseconds')}ms · a dead address says {str(missing.get('detail'))[:40]!r}")

print("── what the deployment is keeping ──")
usage = data_of(call("GET", "/admin/storage", expect=200)) or {}
checks += 1
if not usage.get("databaseBytes"):
    failures.append("the storage screen cannot say how big the database is")
tables = {row.get("name"): row for row in usage.get("tables") or []}
checks += 1
if "presentations" not in tables or "slides" not in tables:
    failures.append(f"the storage screen does not account for the things that grow: {sorted(tables)}")
checks += 1
if any(int(row.get("bytes") or 0) <= 0 for row in tables.values()):
    failures.append("a table is reported as taking no room at all")
call("GET", "/admin/storage", expect=403,
     headers={"X-Ptium-Dev-Secret": SECRET, "Authorization": f"Bearer dev:plain-{RUN}@ptium.local:user"},
     note="what a deployment keeps is the administrator's")
print(f"   database {int(usage.get('databaseBytes', 0))/1024/1024:,.1f} MB across {len(tables)} tables")

# An operator who can see that a deck has waited twenty minutes should be able
# to do something about it. A deck belongs to its owner, and an administrator
# could not see one — let alone push it through or stop it.
print("── the queue an operator can act on ──")
queued_owner = {"X-Ptium-Dev-Secret": SECRET, "Authorization": f"Bearer dev:queue-owner-{RUN}@ptium.local"}
# A big deck on purpose: what this section is about is a stop that lands while a
# worker is holding the deck, and against a four-slide deck the offline writer
# finishes first often enough that the interesting branch is decided by luck.
someone_elses = data_of(call("POST", "/presentations", {"title": f"큐 {RUN}", "prompt": "점검",
                                                        "language": "ko", "slideCount": 30},
                            expect=201, headers=queued_owner)) or {}
call("POST", f"/presentations/{someone_elses['id']}/generate", {}, expect=[200, 202], headers=queued_owner)
waiting = data_of(call("GET", "/admin/generations", expect=200)) or []
mine = [row for row in waiting if row.get("id") == someone_elses["id"]]
checks += 1
if not mine:
    failures.append("a deck waiting in the queue is not in the queue an administrator reads")
else:
    checks += 1
    if not mine[0].get("ownerEmail"):
        failures.append("the queue names its owner by id alone, which nobody knows by sight")
    # Stopping one gives a reason, and the reason is what the author reads —
    # including after the worker that was mid-flight finishes.
    # The answer to the stop is the stopped deck, so the reason is read from it
    # rather than by polling: a poll races the worker and reports the race.
    #
    # The offline writer finishes a deck in about a second, so a stop can
    # honestly arrive too late — and then it must say so rather than answer as
    # though it had stopped something. Both endings are checked; neither is a
    # deck quietly left running.
    status, answered = call("POST", f"/admin/generations/{someone_elses['id']}/cancel",
                            {"reason": f"점검 중단 {RUN}"}, expect=[200, 409])
    if status == 200:
        stopped = data_of((status, answered)) or {}
        checks += 1
        if stopped.get("errorMessage") != f"점검 중단 {RUN}":
            failures.append(f"stopping a deck answers with {stopped.get('errorMessage')!r} rather than the reason given")
        checks += 1
        if stopped.get("status") != "failed":
            failures.append(f"a stopped deck reads as {stopped.get('status')!r}")
        state = data_of(call("GET", f"/presentations/{someone_elses['id']}", expect=200, headers=queued_owner)) or {}
        checks += 1
        if state.get("errorMessage") != f"점검 중단 {RUN}":
            failures.append(f"the author is told {state.get('errorMessage')!r} rather than why it was stopped")
        # The worker was mid-flight when the stop landed; what it writes when it
        # finishes must not replace what the operator said. This belongs to the
        # branch that actually stopped something: after a stop refused as too
        # late there is no reason to overwrite, and asking for one here reported
        # the refusal as a lost reason.
        time.sleep(4)
        state = data_of(call("GET", f"/presentations/{someone_elses['id']}", expect=200, headers=queued_owner)) or {}
        checks += 1
        if state.get("errorMessage") != f"점검 중단 {RUN}":
            failures.append(f"a worker finishing later overwrote the reason an operator gave: "
                            f"{state.get('status')!r} {state.get('errorMessage')!r}")
    else:
        # It finished first. The deck is whole, and the operator was told.
        state = data_of(call("GET", f"/presentations/{someone_elses['id']}", expect=200, headers=queued_owner)) or {}
        checks += 1
        if state.get("status") not in ("completed", "failed"):
            failures.append(f"a stop was refused as too late, but the deck reads as {state.get('status')!r}")
        checks += 1
        if state.get("errorMessage") == f"점검 중단 {RUN}":
            failures.append("a stop was refused and stopped the deck anyway")
    # And it can be pushed back through.
    call("POST", f"/admin/generations/{someone_elses['id']}/requeue", {}, expect=200)
    state = data_of(call("GET", f"/presentations/{someone_elses['id']}", expect=200, headers=queued_owner)) or {}
    checks += 1
    if state.get("status") not in ("queued", "generating", "completed"):
        failures.append(f"a requeued deck reads as {state.get('status')!r}")
    # A deck pushed back into the queue is not the previous attempt's to fail.
    # A worker still holding it reports a moment later, and the fresh attempt
    # died before anybody picked it up.
    state = data_of(call("GET", f"/presentations/{someone_elses['id']}", expect=200, headers=queued_owner)) or {}
    checks += 1
    if state.get("status") == "failed" and "다시 시도" in str(state.get("errorMessage")):
        failures.append("a deck pushed back into the queue was failed by the attempt before it")
    print(f"   one deck: seen, stopped with a reason the author reads, and pushed back")
without_the_role = {"X-Ptium-Dev-Secret": SECRET, "Authorization": f"Bearer dev:plain-{RUN}@ptium.local:user"}
call("GET", "/admin/generations", expect=403, headers=without_the_role,
     note="the queue is the administrator's")
call("POST", f"/admin/generations/{someone_elses['id']}/cancel", {}, expect=403, headers=without_the_role,
     note="stopping someone's deck is the administrator's")

# Everything in this server writes an audit record. Until the trail had a
# reader, "who turned the provider on" and "who deleted that deck" were
# questions with a table behind them and no door.
print("── the trail says who did what ──")
audited = data_of(call("POST", "/presentations", {"title": f"감사 {RUN}", "prompt": "점검",
                                                  "language": "ko"}, expect=201)) or {}
if audited.get("id"):
    call("PUT", f"/presentations/{audited['id']}/source", {"source": "# 한 장\n- 요점\n"}, expect=200)
    call("DELETE", f"/presentations/{audited['id']}", expect=204)
    trail = call("GET", f"/admin/audit?targetId={audited['id']}&days=1", expect=200)[1] or {}
    rows = trail.get("data") if isinstance(trail, dict) else None
    rows = rows if rows is not None else (trail.get("items") if isinstance(trail, dict) else []) or []
    said = [row.get("action") for row in rows]
    checks += 1
    if "presentation.create" not in said or not any(str(a).startswith("presentation.") and "trash" in str(a) or
                                                    str(a) == "presentation.delete" for a in said):
        failures.append(f"the trail of one deck does not hold what was done to it: {said}")
    checks += 1
    if rows and rows[0].get("action") == "presentation.create":
        failures.append("the trail reads oldest first; an operator opens it to see what just happened")
    checks += 1
    if rows and not (rows[0].get("actorEmail") or rows[0].get("actorName")):
        failures.append("the trail names its actor by id alone, which nobody knows by sight")
    actions = data_of(call("GET", "/admin/audit/actions?days=1", expect=200)) or []
    checks += 1
    if not any(row.get("action") == "presentation.create" for row in actions):
        failures.append("the actions a filter can offer do not include what this deployment writes")
    # An auditor asks for the trail as a file, with the filters on screen.
    request = urllib.request.Request(f"{BASE}/admin/audit?format=csv&days=1&action=presentation.create",
                                     headers={"X-Ptium-Dev-Secret": SECRET})
    with urllib.request.urlopen(request, timeout=180) as response:
        exported = response.read()
        disposition = response.headers.get("Content-Disposition") or ""
    checks += 1
    if "attachment" not in disposition or not disposition.endswith('.csv"'):
        failures.append(f"the trail as a file is not offered as one: {disposition!r}")
    checks += 1
    if exported[:3] != b"\xef\xbb\xbf":
        failures.append("the exported trail has no byte order mark; a spreadsheet opens Korean as mojibake")
    exported_rows = list(csv.reader(io.StringIO(exported.decode("utf-8-sig"))))
    checks += 1
    if len(exported_rows) < 2 or exported_rows[0][0] != "시각":
        failures.append(f"the exported trail does not read as a table: {exported_rows[:1]}")
    # The filter travels with the file, and a family name stops at the dot.
    checks += 1
    if any(row[1] != "presentation.create" for row in exported_rows[1:]):
        failures.append("the exported trail carries entries the filter excluded: "
                        f"{sorted({row[1] for row in exported_rows[1:]})}")
    call("GET", "/admin/audit?format=csv", expect=403,
         headers={"X-Ptium-Dev-Secret": SECRET, "Authorization": f"Bearer dev:plain-{RUN}@ptium.local:user"},
         note="the trail as a file is the administrator's too")
    print(f"   one deck: {said} · {len(actions)} kinds of entry today · {len(exported_rows) - 1} rows exported")
# And the trail is the administrator's to read.
plain_headers = {"X-Ptium-Dev-Secret": SECRET, "Authorization": f"Bearer dev:plain-{RUN}@ptium.local:user"}
for path in ("/admin/audit", "/admin/audit/actions"):
    request = urllib.request.Request(BASE + path, headers=plain_headers)
    checks += 1
    try:
        with urllib.request.urlopen(request, timeout=60) as response:
            failures.append(f"{path} answered {response.status} to an account without the administrator role")
    except urllib.error.HTTPError as error:
        if error.code not in (401, 403, 404):
            failures.append(f"{path} answered {error.code} to an account without the administrator role")

# What the queue is doing, which a count on its own does not say.
overview = data_of(call("GET", "/admin/overview", expect=200)) or {}
for field in ("oldestQueuedSeconds", "failedLastDay"):
    checks += 1
    if field not in overview:
        failures.append(f"the overview does not say {field}: a queue depth alone cannot tell a busy morning from a dead worker")
print(f"   queue: {overview.get('queuedGenerations')} waiting, oldest {overview.get('oldestQueuedSeconds')}s, "
      f"{overview.get('failedLastDay')} failed in a day")

# The promise is that the template is the design; the other side of it is that a
# template decides what a deck can be. A design whose only content layout has no
# body region turns every component into a paragraph, and nothing said so until
# the decks came out.
print("── what a template will do to a deck ──")
# "not a usable PowerPoint file" is true of everything rejected. A person
# holding a 97-2003 deck or a file their company's document security wrapped can
# fix it in a minute — if somebody says which of the two it is.
refused = call("POST", "/templates", files={
    "file": ("old.pptx", b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1rest of an OLE file",
             "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    "name": (None, f"옛 형식 {RUN}", None)}, expect=422)
checks += 1
if ".pptx" not in json.dumps(refused[1], ensure_ascii=False):
    failures.append(f"a 97-2003 deck is refused without saying what to do: {refused[1]}")
locked = call("POST", "/templates", files={
    "file": ("locked.pptx", b"SCDSA002\x00\x01locked bytes",
             "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    "name": (None, f"보안 문서 {RUN}", None)}, expect=422)
checks += 1
if "문서보안" not in json.dumps(locked[1], ensure_ascii=False):
    failures.append(f"a file locked by document security is refused without naming it: {locked[1]}")

# A shipped design, not whatever happens to sort first: an account's own upload
# is allowed to be poor, and this check is about what we ship.
shipped = next((t for t in templates if t.get("scope") == "shared"), first)
report = data_of(call("GET", f"/templates/{shipped['id']}/health", expect=200)) or {}
checks += 1
if report.get("slides", 0) < 5:
    failures.append(f"the template report compiled {report.get('slides')} slides of the probe deck")
checks += 1
if not all(report.get("components", {}).get(kind) for kind in ("steps", "kpi", "shareBar", "table", "image")):
    failures.append(f"a shipped template cannot draw the components a brief produces: {report.get('components')}")
checks += 1
if report.get("defects"):
    failures.append(f"a shipped template draws the probe deck with {report.get('defects')} defect(s): "
                    f"{report.get('defectDetails')}")
for role in ("cover", "content", "closing"):
    checks += 1
    if not report.get("roles", {}).get(role):
        failures.append(f"a shipped template has no {role} layout")
print(f"   {report.get('name')}: {report.get('layouts')} layouts, "
      f"{sum(1 for drawn in (report.get('components') or {}).values() if drawn)}/5 components drawn, "
      f"{len(report.get('warnings') or [])} thing(s) it had to say")

# What a slide holds decides what it is. A cover and a closing page hold a name
# and a line or two; a table put on one is written out as text into the room
# kept for that line, and a real weekly report drawn on the title layout every
# week came in as nine lines of text in room for two, on every slide.
# Korean decks number their sections, and every one of those counters is also
# the first syllable of an ordinary word — 개(요), 배(경), 원(인). Read as a
# figure, "1. 개요" is a claim the deck has to cite: the quality panel asked the
# author for the source of "1. 개" on every numbered heading of every deck
# carried in from PowerPoint.
# A table stops at the bottom of its region and caps its columns. What is past
# that is on no slide, and nothing said so: a twelve-row table drew eight rows
# and the panel called the deck clean.
# A report's table is longer than a slide holds. Cutting it at the eighth row is
# how a twelve-row table arrived as eight rows with the rest on no slide and
# nobody told; it continues on the next slide instead, header and all.
print("── a document's long table continues ──")
grid = "구분,내용\n" + "\n".join(f"항목{i},값{i}" for i in range(1, 13))
carried = data_of(call("POST", "/presentations/import",
                       files={"file": (f"실적-{RUN}.csv", grid, "text/csv")}, expect=201)) or {}
carried_id = ((carried.get("presentation") or {}).get("id")) or ""
if carried_id:
    written = (data_of(call("GET", f"/presentations/{carried_id}/source", expect=200)) or {}).get("source", "")
    checks += 1
    missing = [f"항목{i}" for i in range(1, 13) if f"항목{i}" not in written]
    if missing:
        failures.append(f"a twelve-row document table lost {missing}")
    looked = data_of(call("GET", f"/presentations/{carried_id}/inspect", expect=200)) or {}
    checks += 1
    if looked.get("defects"):
        failures.append(f"a document's long table draws with defects: "
                        f"{[f.get('detail') for f in (looked.get('findings') or []) if not f.get('advisory')][:2]}")
    print(f"   rows carried: {12 - len(missing)}/12 · slides: {written.count(chr(35) + ' ')}")
    call("DELETE", f"/presentations/{carried_id}", expect=204)

# A figure gives way to its card down to the size of its label and no further.
# A live model asked for "인력" answered "12 staff redeployed", which is not a
# figure: drawn whole it went through the tile beside it.
print("── a figure too long for its card ──")
carded = data_of(call("POST", "/presentations", {"title": f"지표 {RUN}", "prompt": "점검",
                                                 "language": "ko"}, expect=201))
call("PUT", f"/presentations/{carded['id']}/source",
     {"source": "# 표지\n> 오늘\n\n# 현황과 계획\n@comparison\n> 지금 | 도입 후\n> 수작업 | 자동화\n"
                "::kpi 비교\n- 처리량 | +34%\n- 오류율 | 0.3%\n- 인력 | 12 staff redeployed\n::\n"}, expect=200)
looked = data_of(call("GET", f"/presentations/{carded['id']}/inspect", expect=200)) or {}
checks += 1
if looked.get("defects"):
    failures.append(f"a long figure is drawn over its neighbour: "
                    f"{[f.get('detail') for f in (looked.get('findings') or []) if not f.get('advisory')][:2]}")
checks += 1
if not any(f.get("kind") == "trimmed" for f in (looked.get("findings") or [])):
    failures.append("a figure cut to fit its card said nothing about it")
# An ordinary row of figures is left alone.
call("PUT", f"/presentations/{carded['id']}/source",
     {"source": "# 표지\n> 오늘\n\n# 지표\n::kpi\n- 전환 대상 | 42개\n- 절감 | 18억\n::\n"}, expect=200)
looked = data_of(call("GET", f"/presentations/{carded['id']}/inspect", expect=200)) or {}
checks += 1
if any(f.get("kind") == "trimmed" for f in (looked.get("findings") or [])):
    failures.append("an ordinary row of figures was reported as cut")
print("   long figure: cut and reported · ordinary figures: untouched")
call("DELETE", f"/presentations/{carded['id']}", expect=204)

print("── a table says what it could not draw ──")
wide_table = data_of(call("POST", "/presentations", {"title": f"긴 표 {RUN}", "prompt": "점검",
                                                     "language": "ko"}, expect=201))
rows = "\n".join(f"- 항목{i} | 값{i} | 비고{i}" for i in range(1, 13))
call("PUT", f"/presentations/{wide_table['id']}/source",
     {"source": f"# 표지\n> 오늘\n\n# 열두 줄 표\n::table\n- 구분 | 값 | 비고\n{rows}\n::\n"}, expect=200)
looked = data_of(call("GET", f"/presentations/{wide_table['id']}/inspect", expect=200)) or {}
told = [f.get("detail") for f in (looked.get("findings") or []) if f.get("kind") == "trimmed"]
checks += 1
if not any("rows" in (said or "") for said in told):
    failures.append(f"a twelve-row table drew eight and said {told}")
checks += 1
if looked.get("defects"):
    failures.append(f"a long table is drawn with defects: {looked.get('findings')}")
# And what the file carries is what the preview draws.
status, drawn = call("GET", f"/presentations/{wide_table['id']}/export?format=pptx", raw=True, expect=200)
carried = zipfile.ZipFile(io.BytesIO(drawn)).read("ppt/slides/slide2.xml").decode("utf-8", "replace")
checks += 1
if "항목9" in carried:
    failures.append("the exported file carries a row the slide does not draw")
print(f"   long table: {told}")
call("DELETE", f"/presentations/{wide_table['id']}", expect=204)

# A review is a conversation. A reviewer says the number on slide four is out of
# date; until now the author had nowhere to answer — resolving is a state, and
# "고쳤습니다" is a sentence. An answer beside the remark reads as a second
# remark, so it hangs under it.
print("── a review answers back ──")
reviewed = data_of(call("POST", "/presentations", {"title": f"검토 {RUN}", "prompt": "점검",
                                                   "language": "ko"}, expect=201))
call("PUT", f"/presentations/{reviewed['id']}/source",
     {"source": "# 실적\n- 매출 1,240억\n\n# 계획\n- 요점\n"}, expect=200)
held = data_of(call("GET", f"/presentations/{reviewed['id']}", expect=200)) or {}
first_slide = ((held.get("slides") or [{}])[0]).get("id", "")
opened = data_of(call("POST", f"/presentations/{reviewed['id']}/shares", {}, expect=201)) or {}
review_token = (opened.get("url") or "").rstrip("/").split("/")[-1]
remark = data_of(call("POST", f"/shared/{review_token}/comments",
                      {"slideId": first_slide, "author": "검토자",
                       "body": "이 숫자는 지난 분기 것입니다"}, expect=201, headers={})) or {}
said = data_of(call("POST", f"/presentations/{reviewed['id']}/comments",
                    {"parentId": remark.get("id"), "body": "확인했습니다. 방금 고쳤습니다"}, expect=201)) or {}
checks += 1
if said.get("parentId") != remark.get("id"):
    failures.append(f"an answer was left beside the remark rather than under it: {said!r}")
checks += 1
if said.get("slideId") != first_slide:
    failures.append(f"an answer landed on {said.get('slideId')!r} rather than the remark's slide")
# Answering an answer joins the thread rather than starting a branch.
deeper = data_of(call("POST", f"/shared/{review_token}/comments",
                      {"parentId": said.get("id"), "author": "검토자", "body": "고맙습니다"},
                      expect=201, headers={})) or {}
checks += 1
if deeper.get("parentId") != remark.get("id"):
    failures.append(f"a reply to a reply started a branch: {deeper.get('parentId')!r}")
# Resolving settles the remark and everything under it.
call("POST", f"/presentations/{reviewed['id']}/comments/{remark.get('id')}/resolve",
     {"resolved": True}, expect=[200, 204])
thread = data_of(call("GET", f"/presentations/{reviewed['id']}/comments", expect=200)) or []
checks += 1
if not thread or any(not row.get("resolvedAt") for row in thread):
    failures.append(f"resolving a remark left part of its thread open: {thread!r}")
checks += 1
if call("POST", f"/presentations/{reviewed['id']}/comments",
        {"parentId": "00000000-0000-0000-0000-000000000000", "body": "x"}, expect=422)[0] != 422:
    failures.append("a reply to a comment on no deck was accepted")
# What is waiting on a deck belongs beside its name: somebody with a dozen decks
# had no way to see which two were waiting on them without opening each one.
listed = [row for row in (data_of(call("GET", "/presentations?limit=50", expect=200)) or [])
          if row.get("id") == reviewed["id"]]
checks += 1
if not listed:
    failures.append("a deck with remarks on it is missing from the listing")
elif listed[0].get("openComments"):
    failures.append(f"a settled thread is still counted as waiting: {listed[0].get('openComments')}")
# Another remark, and the count says so — replies are not counted, a thread is
# one thing to deal with.
call("POST", f"/shared/{review_token}/comments",
     {"slideId": first_slide, "author": "검토자", "body": "제목도 봐 주세요"}, expect=201, headers={})
again = [row for row in (data_of(call("GET", "/presentations?limit=50", expect=200)) or [])
         if row.get("id") == reviewed["id"]]
checks += 1
if (again[0].get("openComments") if again else 0) != 1:
    failures.append(f"one remark waiting was counted as {again[0].get('openComments') if again else None}")
print(f"   thread: {len(thread)} messages, all settled together · waiting on the list: {again[0].get('openComments') if again else '-'}")
call("DELETE", f"/presentations/{reviewed['id']}", expect=204)

print("── a numbered heading is not a claim ──")
counted = data_of(call("POST", "/presentations", {"title": f"번호 제목 {RUN}", "prompt": "점검",
                                                  "language": "ko"}, expect=201))
call("PUT", f"/presentations/{counted['id']}/source",
     {"source": "# 표지\n> 오늘\n\n# 1. 개요\n- 상담 품질 관리의 한계\n- 실시간 감정 분석의 필요성\n"
                "\n# 2. 배경\n- 현황을 정리합니다\n- 문제를 짚습니다\n"}, expect=200)
looked = data_of(call("GET", f"/presentations/{counted['id']}/inspect", expect=200)) or {}
cited = [f for f in (looked.get("findings") or []) if f.get("kind") == "source"]
checks += 1
if cited:
    failures.append(f"a numbered heading is read as a figure to cite: {[f.get('detail') for f in cited]}")
# And a real figure is still asked about.
call("PUT", f"/presentations/{counted['id']}/source",
     {"source": "# 표지\n> 오늘\n\n# 개요\n- 매출 1,240억으로 늘었습니다\n- 이익률 9.4%\n"}, expect=200)
looked = data_of(call("GET", f"/presentations/{counted['id']}/inspect", expect=200)) or {}
cited = [f for f in (looked.get("findings") or []) if f.get("kind") == "source"]
checks += 1
if not cited:
    failures.append("a slide of unsourced figures is no longer asked about")
print(f"   numbered headings: clean · figures: {[f.get('detail') for f in cited]}")
call("DELETE", f"/presentations/{counted['id']}", expect=204)

print("── an end page holds what an end page holds ──")
shapely = data_of(call("POST", "/presentations", {"title": f"끝장 {RUN}", "prompt": "점검",
                                                  "language": "ko"}, expect=201))
call("PUT", f"/presentations/{shapely['id']}/source",
     {"source": "# 표지\n> 오늘 결정할 것\n\n# 현황\n- 요점\n\n# 기대효과\n"
                "- 확장성\n- 서비스 범위 확장\n- 다양한 산업군 적용\n- 멀티채널 지원\n- 기능 확장\n- 다국어 지원\n"}, expect=200)
outline = data_of(call("GET", f"/presentations/{shapely['id']}", expect=200)) or {}
last = (outline.get("slides") or [{}])[-1]
checks += 1
if last.get("layout") == "closing":
    failures.append("a last slide of six points was drawn as a closing page")
inspected = data_of(call("GET", f"/presentations/{shapely['id']}/inspect", expect=200)) or {}
checks += 1
if inspected.get("defects"):
    failures.append(f"a deck ending in points has defects: {inspected.get('findings')}")

# A table's rows are a fixed height, and a cell that wraps past its row used to
# be drawn over the row underneath.
call("PUT", f"/presentations/{shapely['id']}/source",
     {"source": "# 표지\n> 오늘\n\n# 플랫폼 활용 방안\n@picture\n::table\n"
                "- 플랫폼 | 활용 방안 | 비고\n"
                "- E-GENE | LCNC 이용한 AI 서비스 개발 | 엔티티 / 폼 / 리스트 / 릴레이션 등 개발\n"
                "- CHEETAH | AI 모델 Proxy | GPU A100\n"
                "- Ollama | Local LLM AI 모델 관리 | AI Model Serving\n::\n"}, expect=200)
inspected = data_of(call("GET", f"/presentations/{shapely['id']}/inspect", expect=200)) or {}
checks += 1
if inspected.get("defects"):
    failures.append(f"a table in a narrow region is drawn with {inspected.get('defects')} defect(s): "
                    f"{[f.get('detail') for f in (inspected.get('findings') or []) if not f.get('advisory')][:2]}")
print(f"   end page: {last.get('layout')!r} · table in a picture layout: {inspected.get('defects')} defect(s)")
call("DELETE", f"/presentations/{shapely['id']}", expect=204)

print("── a deck moved to another template ──")
designs = [t for t in (data_of(call("GET", "/templates?kind=builtin&limit=100", expect=200)) or []) if t.get("id")]
# The destination is a design whose body region is short — the editorial family
# — because that is where carrying the old design's fitting across shows. Two
# designs of the same proportions draw the same deck the same way whether it was
# moved or compiled, and prove nothing.
tight = next((t for t in designs if "Editorial" in (t.get("name") or "")), None)
if tight and len(designs) >= 2:
    designs = [next(t for t in designs if t["id"] != tight["id"]), tight]
if len(designs) >= 2:
    moved = data_of(call("POST", "/presentations", {"title": f"디자인 교체 {RUN}", "prompt": "점검",
                                                    "templateId": designs[0]["id"]}, expect=201)) or {}
    # A deck with components in it, because that is what a template has to hold:
    # a deck of two plain slides moves cleanly into anything.
    moved_source = ("# 표지입니다\n@cover\n> 부제목\n\n"
                    "# 이행 순서\n::steps\n- 준비 | 범위 · 조직 · 예산을 확정\n"
                    "- 이행 | 단계별로 적용하고 완료 조건을 확인\n- 안정화 | 운영 이관과 점검 기준 확정\n::\n\n"
                    "# 기대 효과\n::kpi\n- 전환 시스템 | 42개\n- 절감 | 18억\n- 복구 시간 | 30분\n::\n\n"
                    "# 채널별 비중\n::share\n- 직판 | 46%\n- 대리점 | 33%\n- 온라인 | 21%\n::\n")
    call("PUT", f"/presentations/{moved.get('id')}/source", {"source": moved_source}, expect=200)
    status, drawn_before = call("GET", f"/presentations/{moved.get('id')}/preview.svg?slide=1&width=600", raw=True, expect=200)
    call("PATCH", f"/presentations/{moved.get('id')}", {"templateId": designs[1]["id"]}, expect=200)
    after_state = data_of(call("GET", f"/presentations/{moved.get('id')}", expect=200)) or {}
    status, drawn_after = call("GET", f"/presentations/{moved.get('id')}/preview.svg?slide=1&width=600", raw=True, expect=200)
    if after_state.get("templateName") != designs[1].get("name"):
        failures.append(f"the deck did not move to {designs[1].get('name')}: {after_state.get('templateName')}")
    if drawn_before == drawn_after:
        failures.append("a deck moved to another template draws exactly as it did before")
    moved_measured = data_of(call("GET", f"/presentations/{moved.get('id')}/inspect", expect=200)) or {}
    if moved_measured.get("defects"):
        failures.append(f"the deck has {moved_measured.get('defects')} defect(s) in its new design")
    status, moved_pptx = call("GET", f"/presentations/{moved.get('id')}/export?format=pptx", raw=True, expect=200)
    if not moved_pptx.startswith(b"PK"):
        failures.append("a deck in its new design does not export")
    # Moving a deck into a template has to leave it as the template would have
    # drawn it. The slides carry the sizes and the cuts decided for the design
    # they were compiled in, so a deck that moved kept the old design's fitting
    # and drew outside its regions in a template that draws it cleanly.
    born = data_of(call("POST", "/presentations", {"title": f"디자인 원본 {RUN}", "prompt": "점검",
                                                   "templateId": designs[1]["id"]}, expect=201)) or {}
    call("PUT", f"/presentations/{born.get('id')}/source", {"source": moved_source}, expect=200)
    born_measured = data_of(call("GET", f"/presentations/{born.get('id')}/inspect", expect=200)) or {}
    checks += 1
    if moved_measured.get("defects") != born_measured.get("defects"):
        failures.append(f"a deck moved to {designs[1].get('name')} has {moved_measured.get('defects')} defect(s) "
                        f"and the same deck compiled into it has {born_measured.get('defects')}")
    # The same deck and the same design draw the same slide, however it got
    # there. Comparing the drawings rather than the defect count catches a deck
    # that kept the old design's sizes in a template where that still fits.
    for position in (2, 3, 4):
        status, moved_drawing = call("GET", f"/presentations/{moved.get('id')}/preview.svg?slide={position}&width=600",
                                     raw=True, expect=200)
        status, born_drawing = call("GET", f"/presentations/{born.get('id')}/preview.svg?slide={position}&width=600",
                                    raw=True, expect=200)
        checks += 1
        if moved_drawing != born_drawing:
            failures.append(f"slide {position} drawn one way when the deck moved to {designs[1].get('name')} "
                            f"and another way when it was compiled into it")
    call("DELETE", f"/presentations/{born.get('id')}", expect=204)
    # The workspace sends the slides it has and the new design in one request,
    # and when the deck is dirty those slides hold edits the stored source has
    # never seen. Compiling the stored source would throw them away.
    edited = data_of(call("POST", "/presentations", {"title": f"동시 편집 {RUN}", "prompt": "점검",
                                                     "templateId": designs[0]["id"]}, expect=201)) or {}
    call("PUT", f"/presentations/{edited.get('id')}/source",
         {"source": "# 첫 장\n- 원래 요점\n\n# 둘째 장\n- 그대로 둘 요점\n"}, expect=200)
    state = data_of(call("GET", f"/presentations/{edited.get('id')}", expect=200)) or {}
    sending = []
    for index, slide in enumerate(state.get("slides") or []):
        body = slide.get("content")
        if isinstance(body, str):
            body = json.loads(body)
        if index == 0:
            body.setdefault("fields", {})["title"] = [{"text": "편집한 제목"}]
        sending.append({"title": "편집한 제목" if index == 0 else slide.get("title"), "content": body,
                        "layout": slide.get("layout"), "layoutId": slide.get("layoutId"),
                        "speakerNotes": slide.get("speakerNotes") or ""})
    call("PUT", f"/presentations/{edited.get('id')}",
         {"slides": sending, "templateId": designs[1]["id"], "version": state.get("version")}, expect=200)
    carried = data_of(call("GET", f"/presentations/{edited.get('id')}", expect=200)) or {}
    checks += 1
    if ((carried.get("slides") or [{}])[0]).get("title") != "편집한 제목":
        failures.append("an edit sent with a design change was thrown away by the recompile")
    written = (data_of(call("GET", f"/presentations/{edited.get('id')}/source", expect=200)) or {}).get("source", "")
    checks += 1
    if "편집한 제목" not in written:
        failures.append("an edit sent with a design change is not in the deck's source")
    call("DELETE", f"/presentations/{edited.get('id')}", expect=204)

    # Compiling a deck again is a rewrite of every slide, so what is not written
    # in the source has to be carried across it: an image placed in a slot, an
    # object the author put on the canvas, and the notes.
    carried_over = data_of(call("POST", "/presentations", {"title": f"교체 보존 {RUN}", "prompt": "점검",
                                                           "templateId": designs[0]["id"]}, expect=201)) or {}
    call("PUT", f"/presentations/{carried_over.get('id')}/source",
         {"source": f"# 그림 장\n@picture\n::image logo-{RUN}.png | 현장\n!notes 노트는 남는다\n"}, expect=200)
    holding = data_of(call("GET", f"/presentations/{carried_over.get('id')}", expect=200)) or {}
    kept = (holding.get("slides") or [{}])[0]
    kept_content = kept.get("content")
    if isinstance(kept_content, str):
        kept_content = json.loads(kept_content)
    kept_content["elements"] = [{"id": "el-1", "kind": "text", "text": "직접 놓은 글",
                                 "x": 0.1, "y": 0.1, "width": 0.3, "height": 0.1, "zIndex": 1}]
    call("PUT", f"/presentations/{carried_over.get('id')}",
         {"slides": [{"title": kept.get("title"), "content": kept_content, "layout": kept.get("layout"),
                      "layoutId": kept.get("layoutId"), "speakerNotes": kept.get("speakerNotes") or ""}],
          "version": holding.get("version")}, expect=200)
    holding = data_of(call("GET", f"/presentations/{carried_over.get('id')}", expect=200)) or {}
    call("PUT", f"/presentations/{carried_over.get('id')}",
         {"templateId": designs[1]["id"], "version": holding.get("version")}, expect=200)
    moved_state = data_of(call("GET", f"/presentations/{carried_over.get('id')}", expect=200)) or {}
    surviving = (moved_state.get("slides") or [{}])[0]
    surviving_content = surviving.get("content")
    if isinstance(surviving_content, str):
        surviving_content = json.loads(surviving_content)
    checks += 1
    if not (surviving_content.get("images") or {}):
        failures.append("a deck moved to another template lost the image placed on it")
    checks += 1
    if not (surviving_content.get("elements") or []):
        failures.append("a deck moved to another template lost what the author put on the canvas")
    checks += 1
    if "노트는 남는다" not in (surviving.get("speakerNotes") or ""):
        failures.append("a deck moved to another template lost its speaker notes")
    print("   image, canvas objects and notes carried across the move")
    call("DELETE", f"/presentations/{carried_over.get('id')}", expect=204)
    print(f"   {designs[0].get('name')} → {designs[1].get('name')}: drew differently, "
          f"{moved_measured.get('defects')} defect(s), same as compiled into it, edits kept")
    call("DELETE", f"/presentations/{moved.get('id')}", expect=204)

print("── the numbers a brief gives, drawn ──")
figures = data_of(call("POST", "/presentations", {"title": f"숫자 브리프 {RUN}",
    "prompt": "지역별 처리 건수 보고. 서울 1,200건, 부산 860건, 대구 540건, 광주 410건. "
              "채널 비중은 직판 46%, 대리점 33%, 온라인 21%.",
    "requestedSlideCount": 8, "language": "ko"}, expect=201)) or {}
call("POST", f"/presentations/{figures.get('id')}/generate", {}, expect=[200, 202])
for _ in range(90):
    time.sleep(2)
    drawn_state = data_of(call("GET", f"/presentations/{figures.get('id')}", expect=200)) or {}
    if drawn_state.get("status") in ("completed", "failed"):
        break
drawn_source = (data_of(call("GET", f"/presentations/{figures.get('id')}/source", expect=200)) or {}).get("source", "")
kinds = re.findall(r"^::(\w+)", drawn_source, re.M)
print("   components:", kinds)
if not [kind for kind in kinds if kind in ("columns", "bars", "share", "line", "meter")]:
    failures.append(f"a brief of nothing but numbers drew no chart: {kinds}")
# What a component row is named. Prose may quote the brief; a bar may not — its
# label is all a reader gets.
component_rows = [line for line in drawn_source.split("\n") if line.startswith("- ") and " | " in line]
for row in component_rows:
    label = row[2:].split(" | ", 1)[0]
    if label.startswith("보고.") or label.startswith("비중은") or re.search(r"[0-9]", label):
        failures.append(f"a component row is named by the words around the number: {row!r}")
        break
measured_figures = data_of(call("GET", f"/presentations/{figures.get('id')}/inspect", expect=200)) or {}
if measured_figures.get("defects"):
    failures.append(f"the deck that draws its figures has {measured_figures.get('defects')} defect(s)")
call("DELETE", f"/presentations/{figures.get('id')}", expect=204)

# A deck can only carry a picture somebody already uploaded — an air-gapped
# deployment has no web to fetch one from — so generation is told what this
# account has and places the one a slide is named for.
# A deck-wide rewrite takes what the author asked for, in their own words. It
# needs a provider to run, so what is checked without one is that the words are
# accepted, bounded, and refused for the right reason.
call("POST", f"/presentations/{deck_id}/rewrite", {"instruction": "2장은 요점 3개로 줄여 주세요"}, expect=409,
     note="rewriting needs a provider, and says so")
status, refused_long = call("POST", f"/presentations/{deck_id}/rewrite", {"instruction": "가" * 2100}, expect=422,
                            note="an instruction longer than the model will read is refused")
if (refused_long or {}).get("error", {}).get("code") != "validation_error":
    failures.append(f"an over-long instruction is refused for the wrong reason: {refused_long}")

print("── the pictures an account already has ──")
picture = data_of(call("POST", "/assets", files={"file": (f"현장 자동화 {RUN}.png", png(rgb=(30, 120, 200)), "image/png")},
                       expect=201)) or {}
call("PATCH", f"/assets/{picture.get('id')}", {"name": f"현장 자동화 {RUN}"}, expect=200)
illustrated = data_of(call("POST", "/presentations", {"title": f"현장 자동화 {RUN}",
    "prompt": f"현장 자동화 {RUN} 도입 계획을 임원에게 보고합니다. 투자 12억, 12개월.",
    "requestedSlideCount": 6, "language": "ko"}, expect=201)) or {}
call("POST", f"/presentations/{illustrated.get('id')}/generate", {}, expect=[200, 202])
for _ in range(90):
    time.sleep(2)
    illustrated_state = data_of(call("GET", f"/presentations/{illustrated.get('id')}", expect=200)) or {}
    if illustrated_state.get("status") in ("completed", "failed"):
        break
carrying = [slide for slide in (illustrated_state.get("slides") or []) if (slide.get("content") or {}).get("images")]
print(f"   slides carrying the account's picture: {len(carrying)}")
if not carrying:
    failures.append("a deck about a picture the account has was written without it")
illustrated_measured = data_of(call("GET", f"/presentations/{illustrated.get('id')}/inspect", expect=200)) or {}
if illustrated_measured.get("defects"):
    failures.append(f"the illustrated deck has {illustrated_measured.get('defects')} defect(s)")
call("DELETE", f"/presentations/{illustrated.get('id')}", expect=204)
call("DELETE", f"/assets/{picture.get('id')}", expect=204)

print("── generation without an AI provider (the air-gapped path) ──")
draft = data_of(call("POST", "/presentations", {
    "title": "오프라인 생성", "prompt": "2026년 하반기 클라우드 전환 로드맵과 투자 타당성을 6장으로 정리해줘",
    "requestedSlideCount": 6, "language": "ko", "templateId": first["id"]}, expect=201))
call("POST", f"/presentations/{draft['id']}/generate", {}, expect=[200, 202])
# Without a provider the deck is written locally in a second; with one it is a
# round trip to a model, which on a self-hosted one is minutes. Both are valid
# deployments, so the sweep waits and then says which it saw.
for _ in range(240):
    state = data_of(call("GET", f"/presentations/{draft['id']}", expect=200))
    if state["status"] in ("completed", "failed"):
        break
    time.sleep(1)
print("   status:", state["status"], "slides:", len(state.get("slides") or []))
if state["status"] not in ("completed", "failed") :
    print("   (still with the model after four minutes; the rest of this section is skipped)")
elif state["status"] != "completed":
    failures.append(f"generation ended as {state['status']}: {state.get('errorMessage')}")
elif len(state.get("slides") or []) != 6:
    failures.append(f"asked for 6 slides, generated {len(state.get('slides') or [])}")
if state["status"] == "completed":
    generated = data_of(call("GET", f"/presentations/{draft['id']}/inspect", expect=200))
    print("   generated deck defects:", (generated or {}).get("defects"), "advisories:", (generated or {}).get("advisories"))
    if (generated or {}).get("defects"):
        failures.append(f"the generated deck has defects: {generated.get('findings')}")
    # A deck can be drawn perfectly and headed with half a sentence. The
    # measurement says so now, so the sweep can ask it.
    checks += 1
    unfinished = [f for f in ((generated or {}).get("findings") or []) if f.get("kind") == "unfinished"]
    if unfinished:
        failures.append(f"the generated deck is headed with a fragment: {unfinished}")
    call("DELETE", f"/presentations/{draft['id']}", expect=204)

# A brief that asks for a section is not a brief about the asking. "마지막에
# 리스크 장을 꼭 넣어주세요" used to arrive as a slide titled with the sentence
# itself, twice over.
asked_deck = data_of(call("POST", "/presentations", {
    "title": "", "prompt": "물류센터 자동화 도입을 임원에게 보고합니다. 마지막에 리스크 장을 꼭 넣어주세요.",
    "requestedSlideCount": 8, "language": "ko", "templateId": first["id"]}, expect=201))
call("POST", f"/presentations/{asked_deck['id']}/generate", {}, expect=[200, 202])
for _ in range(240):
    asked_state = data_of(call("GET", f"/presentations/{asked_deck['id']}", expect=200))
    if asked_state["status"] in ("completed", "failed"):
        break
    time.sleep(1)
if asked_state["status"] == "completed":
    asked_source = (data_of(call("GET", f"/presentations/{asked_deck['id']}/source", expect=200)) or {}).get("source", "")
    titles = [line[2:] for line in asked_source.splitlines() if line.startswith("# ")]
    checks += 1
    if not any("리스크" in title for title in titles):
        failures.append(f"the section the brief asked for is not in the deck: {titles}")
    checks += 1
    said = [title for title in titles if "넣어" in title or "슬라이드" in title or "마지막에" in title]
    if said:
        failures.append(f"a slide is titled with the request itself: {said}")
    measured = data_of(call("GET", f"/presentations/{asked_deck['id']}/inspect", expect=200)) or {}
    checks += 1
    fragments = [f for f in (measured.get("findings") or []) if f.get("kind") == "unfinished"]
    if fragments:
        failures.append(f"a deck written from a brief is headed with a fragment: {fragments}")
    print("   asked-for section:", " · ".join(titles)[:70])
call("DELETE", f"/presentations/{asked_deck['id']}", expect=204)

# Every other number on the admin overview counts decks that have not been
# deleted. Without this one, an operator cannot see that a deployment is holding
# thousands of decks nobody wanted — nothing empties the trash on its own.
overview = data_of(call("GET", "/admin/overview", expect=[200, 403]))
if overview is not None and isinstance(overview, dict) and "presentations" in overview:
    checks += 1
    if "deletedDecks" not in overview:
        failures.append(f"the admin overview does not say what the trash holds: {sorted(overview)}")
    else:
        print("   trash:", overview.get("deletedDecks"), "deck(s) · live:", overview.get("presentations"))

# What the API answers with has to be what the schema says it answers with: the
# overview's envelope is closed (additionalProperties: false), so a client that
# validates the response refuses it over a field nobody documented.
if overview is not None and isinstance(overview, dict) and "presentations" in overview:
    documented = set()
    try:
        import re as _re
        text = open(os.path.join(os.path.dirname(__file__), "..", "..", "api", "openapi.yaml"),
                    encoding="utf-8").read()
        block = text.split("    AdminOverview:", 1)[1].split("\n    AdminOverviewEnvelope:", 1)[0]
        documented = set(_re.findall(r"\n        (\w+):", block))
    except Exception as reason:
        failures.append(f"the sweep could not read the overview schema: {reason}")
    checks += 1
    undocumented = sorted(set(overview) - documented)
    if documented and undocumented:
        failures.append(f"the admin overview answers with fields the schema does not list: {undocumented}")

# A colour somebody chose has to reach the drawing. Every slide carried the
# author's brand colour, the profile screen said it would be used, and nothing
# read it.
was = data_of(call("GET", "/profile", expect=200)) or {}
call("PATCH", "/profile", {"preferences": {**(was.get("preferences") or {}), "brandColor": "#0F62FE"}}, expect=200)
branded = data_of(call("POST", "/presentations", {"title": f"브랜드 색 {RUN}",
                                                  "prompt": "클라우드 전환 로드맵과 투자 타당성을 정리해 주세요",
                                                  "requestedSlideCount": 6, "language": "ko"}, expect=201)) or {}
if branded.get("id"):
    call("POST", f"/presentations/{branded['id']}/generate", {}, expect=[200, 202])
    for _ in range(240):
        state = data_of(call("GET", f"/presentations/{branded['id']}", expect=200)) or {}
        if state.get("status") in ("completed", "failed"):
            break
        time.sleep(1)
    if state.get("status") == "completed":
        painted = []
        for position in range(1, len(state.get("slides") or []) + 1):
            _, drawing = call("GET", f"/presentations/{branded['id']}/preview.svg?slide={position}&width=900",
                              raw=True, expect=200)
            if b"0F62FE" in (drawing or b"").upper():
                painted.append(position)
        checks += 1
        if not painted:
            failures.append("a brand colour was set and no slide was drawn with it")
        print("   slides drawn in the author's colour:", painted)
    call("DELETE", f"/presentations/{branded['id']}?permanent=true", expect=[204, 404])
call("PATCH", "/profile", {"preferences": was.get("preferences") or {}}, expect=200)

print("── api keys and mcp ──")
made = data_of(call("POST", "/api-keys", {"name": f"e2e {RUN}", "scopes": ["presentations:read"]}, expect=201)) or {}
call("GET", "/api-keys", expect=200)
# The secret comes back once, under its own name; the key itself is beside it.
# Reading it from the wrong field made every check below quietly skip — the
# account collected eighty-nine keys and nothing had used one since.
key = made.get("apiKey") or made.get("api_key") or made
token = made.get("secret") or made.get("key") or made.get("token") or ""
checks += 1
if not token or not key.get("id"):
    failures.append(f"a new api key came back without a secret to use: {sorted(made)}")
if token and key.get("id"):
    request = urllib.request.Request(BASE + "/presentations", headers={"Authorization": "Bearer " + token})
    try:
        with urllib.request.urlopen(request) as response:
            checks += 1
            if response.status != 200:
                failures.append(f"an api key could not read presentations: {response.status}")
    except urllib.error.HTTPError as error:
        failures.append(f"an api key could not read presentations: {error.code} {error.read()[:200]!r}")
    # A read-only key must not write.
    request = urllib.request.Request(BASE + "/presentations", data=b"{}", method="POST",
                                     headers={"Authorization": "Bearer " + token, "Content-Type": "application/json"})
    try:
        urllib.request.urlopen(request)
        failures.append("a read-only api key was allowed to create a deck")
    except urllib.error.HTTPError as error:
        checks += 1
        if error.code != 403:
            failures.append(f"a read-only key writing gave {error.code}, expected 403")
    # What a key may do can be changed without changing the key: the key is
    # written into a machine's configuration and forgotten, and issuing another
    # one to add a permission means going back there.
    offered = data_of(call("GET", "/api-keys/scopes", expect=200)) or []
    checks += 1
    if not {"templates:read", "presentations:write"} <= {scope.get("id") for scope in offered}:
        failures.append(f"the scopes a key may carry are not all offered: {[s.get('id') for s in offered]}")
    request = urllib.request.Request(BASE + "/templates?limit=1", headers={"Authorization": "Bearer " + token})
    try:
        urllib.request.urlopen(request)
        failures.append("a key without templates:read read the templates")
    except urllib.error.HTTPError as error:
        checks += 1
        if error.code != 403:
            failures.append(f"a key without templates:read gave {error.code}, expected 403")
    widened = data_of(call("PATCH", f"/api-keys/{key['id']}",
                           {"scopes": ["presentations:read", "templates:read"]}, expect=200)) or {}
    checks += 1
    if "templates:read" not in (widened.get("scopes") or []):
        failures.append(f"widening a key's scopes answered {widened.get('scopes')!r}")
    try:
        with urllib.request.urlopen(urllib.request.Request(
                BASE + "/templates?limit=1", headers={"Authorization": "Bearer " + token})) as response:
            checks += 1
            if response.status != 200:
                failures.append(f"a widened key still cannot read templates: {response.status}")
    except urllib.error.HTTPError as error:
        failures.append(f"a widened key still cannot read templates: {error.code}")

    # A key cannot carry more than its owner has. The catalogue marks the admin
    # scopes, minting refuses them to an ordinary owner, and so does widening a
    # key that already exists — the third door being the one somebody would try
    # after the first two closed.
    plain_owner = {"X-Ptium-Dev-Secret": SECRET,
                   "Authorization": f"Bearer dev:plain-{RUN}@ptium.local:user",
                   "Content-Type": "application/json"}
    for_plain = data_of(call("GET", "/api-keys/scopes", expect=200, headers=plain_owner)) or []
    checks += 1
    admin_offered = [scope.get("id") for scope in for_plain if str(scope.get("id")).startswith("admin:")]
    if admin_offered:
        failures.append(f"an ordinary owner is offered the admin scopes {admin_offered}")
    for wanted in (["admin:settings"], ["presentations:read", "admin:users"]):
        status, body = call("POST", "/api-keys", {"name": f"escalation {RUN}", "scopes": wanted},
                            expect=422, headers=plain_owner, note=f"an ordinary owner minting {wanted}")
        if status not in (401, 403, 422):
            failures.append(f"an ordinary owner minting {wanted} answered {status}")
    ordinary = data_of(call("POST", "/api-keys", {"name": f"ordinary {RUN}", "scopes": ["presentations:read"]},
                            expect=201, headers=plain_owner)) or {}
    ordinary_key = (ordinary.get("apiKey") or ordinary)
    if ordinary_key.get("id"):
        status, _ = call("PATCH", f"/api-keys/{ordinary_key['id']}",
                         {"scopes": ["presentations:read", "admin:settings"]},
                         expect=422, headers=plain_owner, note="an ordinary owner widening to an admin scope")
        if status not in (401, 403, 422):
            failures.append(f"an ordinary owner widened a key to admin:settings ({status})")
        call("DELETE", f"/api-keys/{ordinary_key['id']}", expect=(200, 204), headers=plain_owner)
    call("PATCH", f"/api-keys/{key['id']}", {"scopes": ["nope:everything"]}, expect=422)
    call("DELETE", f"/api-keys/{key['id']}", expect=[204, 200])
    # Deleting a key revokes it: the row stays so the audit trail still names
    # what used it. What must not stay is a key that still works.
    mine = [row for row in (data_of(call("GET", "/api-keys", expect=200)) or []) if RUN in (row.get("name") or "")]
    checks += 1
    if any(not row.get("revokedAt") for row in mine):
        failures.append(f"this run left a usable api key behind: {[row.get('name') for row in mine]}")

# What a sweep uploads, a sweep takes away. The account's library is what the
# writer offers a deck when it looks for a picture, so a run's leftovers end up
# on somebody's slide — one did, on a deck written by the live model.
print("── clearing what this run uploaded ──")
call("DELETE", f"/assets/{asset_id}", expect=[204, 404])
leftover = [a for a in (data_of(call("GET", "/assets?limit=200", expect=200)) or [])
            if RUN in (a.get("name") or "")]
checks += 1
if leftover:
    failures.append(f"this run left {len(leftover)} image(s) in the library: "
                    f"{[a.get('name') for a in leftover][:3]}")

print("── authentication and errors ──")
request = urllib.request.Request(BASE + "/presentations")
try:
    urllib.request.urlopen(request)
    failures.append("an unauthenticated request was allowed")
except urllib.error.HTTPError as error:
    checks += 1
    if error.code != 401:
        failures.append(f"unauthenticated gave {error.code}, expected 401")
call("GET", "/presentations/00000000-0000-0000-0000-000000000000", expect=404)
call("POST", "/presentations", {"title": ""}, expect=[201, 422])
call("GET", "/nope", expect=404)

print()
print(f"{checks} checks, {len(failures)} failures")
for failure in failures:
    print(" ✗ " + failure)
sys.exit(1 if failures else 0)
